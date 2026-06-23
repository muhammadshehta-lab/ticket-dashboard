import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import streamlit as st
import gspread
from google.oauth2.service_account import Credentials
import json, hashlib, time, pathlib, urllib.parse, re, requests
from datetime import timedelta
import smtplib
from email.message import EmailMessage
import base64
from io import BytesIO
from PIL import Image

# ══════════════════════════════════════════════════════════════════════════════════
#  PAGE CONFIG
# ══════════════════════════════════════════════════════════════════════════════════
st.set_page_config(
    page_title="In-Store Requests Dashboard",
    page_icon="💊",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ══════════════════════════════════════════════════════════════════════════════════
#  PERSISTENCE & PRE-SEEDED USER DATABASE
# ══════════════════════════════════════════════════════════════════════════════════
_DATA_FILE = pathlib.Path(__file__).parent / ".dashboard_data.json"

def _load_store() -> dict:
    default_store = {
        "users": {
            "admin": {"display_name": "Mohammed Shehta", "password_hash": hashlib.sha256("admin123".encode()).hexdigest(), "role": "admin", "agent_name": None},
            "50107": {"display_name": "Ahmed El-Kholy", "password_hash": hashlib.sha256("50107".encode()).hexdigest(), "role": "expert", "agent_name": "Ahmed El-Kholy"},
            "50399": {"display_name": "Ahmed Kadry", "password_hash": hashlib.sha256("50399".encode()).hexdigest(), "role": "expert", "agent_name": "Ahmed Kadry"},
            "50187": {"display_name": "Amr El-Sayed", "password_hash": hashlib.sha256("50187".encode()).hexdigest(), "role": "expert", "agent_name": "Amr El-Sayed"},
            "50461": {"display_name": "Eslam Ramadan", "password_hash": hashlib.sha256("50461".encode()).hexdigest(), "role": "expert", "agent_name": "Eslam Ramadan"},
            "50274": {"display_name": "Mohamed Abdelmageed", "password_hash": hashlib.sha256("50274".encode()).hexdigest(), "role": "expert", "agent_name": "Mohamed Abdelmageed"},
            "50476": {"display_name": "Mohamed Khalifa", "password_hash": hashlib.sha256("50476".encode()).hexdigest(), "role": "expert", "agent_name": "Mohamed Khalifa"},
            "50114": {"display_name": "Yahia Ali Shafei", "password_hash": hashlib.sha256("50114".encode()).hexdigest(), "role": "expert", "agent_name": "Yahia Ali Shafei"}
        },
        "requests":  [], "overrides": {}, "login_logs": [],
    }
    if _DATA_FILE.exists():
        try:
            loaded = json.loads(_DATA_FILE.read_text())
            if "users" in loaded:
                for k, v in loaded["users"].items(): default_store["users"][k] = v
            if "requests" in loaded: default_store["requests"] = loaded["requests"]
            if "overrides" in loaded: default_store["overrides"] = loaded["overrides"]
            if "login_logs" in loaded: default_store["login_logs"] = loaded["login_logs"]
        except Exception: pass
    return default_store

def _save_store(): _DATA_FILE.write_text(json.dumps(st.session_state.store, indent=2))
def _hash(pw: str) -> str: return hashlib.sha256(pw.encode()).hexdigest()

def parse_drive_link(raw_url):
    raw_url = str(raw_url).strip()
    if "drive.google.com" in raw_url:
        try:
            match = re.search(r'/d/([a-zA-Z0-9_-]+)', raw_url)
            if not match: match = re.search(r'id=([a-zA-Z0-9_-]+)', raw_url)
            if match: return f"https://drive.google.com/thumbnail?id={match.group(1)}&sz=w500"
        except: return raw_url
    return raw_url

# ══════════════════════════════════════════════════════════════════════════════════
#  GOOGLE SHEETS SYNC HELPERS (PASSWORD & PROFILE)
# ══════════════════════════════════════════════════════════════════════════════════
def update_sheet_password(uname: str, new_pass: str):
    try:
        scopes = ["https://www.googleapis.com/auth/spreadsheets", "https://www.googleapis.com/auth/drive"]
        if "gspread" in st.secrets and "credentials" in st.secrets["gspread"]:
            creds = Credentials.from_service_account_info(json.loads(st.secrets["gspread"]["credentials"]), scopes=scopes)
            client = gspread.authorize(creds)
            ws = client.open("AlDawaa Tickets Data").worksheet("Users")
            cell = ws.find(str(uname), in_column=1)
            if cell: ws.update_cell(cell.row, 2, str(new_pass))
    except Exception as e: pass

def update_sheet_profile(uname: str, profile_data: dict):
    try:
        scopes = ["https://www.googleapis.com/auth/spreadsheets", "https://www.googleapis.com/auth/drive"]
        if "gspread" in st.secrets and "credentials" in st.secrets["gspread"]:
            creds = Credentials.from_service_account_info(json.loads(st.secrets["gspread"]["credentials"]), scopes=scopes)
            client = gspread.authorize(creds)
            ws = client.open("AlDawaa Tickets Data").worksheet("Users")
            cell = ws.find(str(uname), in_column=1)
            if cell:
                row = cell.row
                ws.update_cell(row, 7, str(profile_data.get('photo', '')))
                ws.update_cell(row, 8, str(profile_data.get('grad_year', '')))
                ws.update_cell(row, 9, str(profile_data.get('join_cc', '')))
                ws.update_cell(row, 10, str(profile_data.get('join_team', '')))
                ws.update_cell(row, 11, str(profile_data.get('bio', '')))
    except Exception as e: pass

# ══════════════════════════════════════════════════════════════════════════════════
#  NOTIFICATIONS (WHATSAPP & EMAIL)
# ══════════════════════════════════════════════════════════════════════════════════
def notify_admin_whatsapp(logged_in_user):
    try:
        if "whatsapp" in st.secrets and "api_key" in st.secrets["whatsapp"]:
            api_key = st.secrets["whatsapp"]["api_key"]
            phone = "+201129217380"
            msg = f"🚨 *System Login Alert*%0AUser: *{logged_in_user}*%0ATime: {time.strftime('%Y-%m-%d %H:%M:%S')}"
            url = f"https://api.callmebot.com/whatsapp.php?phone={phone}&text={msg}&apikey={api_key}"
            requests.get(url, timeout=3)
    except Exception as e: pass

def send_approval_email(to_email, name, uid):
    try:
        if "smtp" in st.secrets and "email" in st.secrets["smtp"] and "password" in st.secrets["smtp"]:
            sender_email = st.secrets["smtp"]["email"]
            sender_password = st.secrets["smtp"]["password"]
            msg = EmailMessage()
            msg['Subject'] = "✅ AlDawaa Dashboard Access Approved"
            msg['From'] = f"Mohammed Shehta <{sender_email}>"
            msg['To'] = to_email
            dashboard_url = "https://aldawaa-requests.streamlit.app" 
            body = f"Dear {name},\n\nYour request to access the AlDawaa In-Store Requests Dashboard has been successfully approved.\n\nHere are your login credentials:\n- Username / ID: {uid}\n- Temporary Password: {uid}\n\n(Note: You will be required to set a new, secure password upon your first login).\n\nYou can access the dashboard via the link below:\n{dashboard_url}\n\nBest regards,\nMohammed Shehta"
            msg.set_content(body)
            with smtplib.SMTP('smtp.gmail.com', 587) as server:
                server.starttls()
                server.login(sender_email, sender_password)
                server.send_message(msg)
            return True
        return False
    except Exception as e: return False

# ══════════════════════════════════════════════════════════════════════════════════
#  CSS  — PROFESSIONAL CORPORATE THEME & PROFILE CARDS
# ══════════════════════════════════════════════════════════════════════════════════
st.markdown("""
<style>
.stApp { background: #f1f5f9 !important; color: #0f172a !important; }
.stApp, p, span, label, input, select { font-size: 1.05rem !important; font-weight: 600 !important; }
h1 { font-size: 2.5rem !important; font-weight: 900 !important; color: #0f172a !important; }
h2 { font-size: 2.15rem !important; font-weight: 900 !important; color: #1e293b !important; margin-top: 1rem !important; }
h3 { font-size: 1.65rem !important; font-weight: 800 !important; color: #1e293b !important; }
h4 { font-size: 1.35rem !important; font-weight: 800 !important; color: #334155 !important; }

[data-testid="stSidebar"] { background: #ffffff !important; border-right: 1px solid #e2e8f0 !important; }
[data-testid="stSidebar"] * { color: #1e293b !important; font-weight: 700 !important; }
[data-testid="stTabs"] [data-baseweb="tab"] { background: transparent; color: #64748b !important; font-weight: 700 !important; font-size: 1.15rem !important; padding-bottom: 8px !important; }
[data-testid="stTabs"] [aria-selected="true"] { color: #2563eb !important; border-bottom: 3px solid #2563eb !important; }

.stTextInput input, .stNumberInput input, .stSelectbox select, .stDateInput input { background: #ffffff !important; border: 1px solid #cbd5e1 !important; color: #0f172a !important; border-radius: 8px !important; font-weight: 600 !important; font-size: 1.05rem !important; }
.stTextInput input:focus { border-color: #2563eb !important; box-shadow: 0 0 0 3px rgba(37,99,235,0.1) !important; }
.stButton > button { background: #2563eb !important; border: 1px solid #1d4ed8 !important; color: #ffffff !important; border-radius: 8px !important; font-weight: 700 !important; font-size: 1.1rem !important; padding: 0.5rem 1.5rem !important; box-shadow: 0 2px 4px rgba(37,99,235,0.2); }
.stButton > button:hover { background: #1d4ed8 !important; border-color: #1e40af !important; transform: translateY(-1px); box-shadow: 0 4px 6px rgba(37,99,235,0.3); }

.kpi-container { position: relative; border-radius: 12px; padding: 1.4rem 1.1rem; text-align: center; min-height: 124px; display: flex; flex-direction: column; justify-content: center; margin-bottom: 1rem; box-shadow: 0 1px 3px rgba(0,0,0,0.05), 0 1px 2px rgba(0,0,0,0.03); border-width: 1px !important; border-style: solid !important; background: #ffffff; }
.kpi-container:hover { transform: translateY(-2px); box-shadow: 0 4px 6px rgba(0,0,0,0.07), 0 2px 4px rgba(0,0,0,0.05); }
.kpi-label { font-size: .8rem !important; letter-spacing: .05em; text-transform: uppercase; margin-bottom: .5rem; font-weight: 700 !important; color: #64748b !important; }
.kpi-value { font-size: 1.9rem !important; font-weight: 800 !important; letter-spacing: -.02em; color: #0f172a !important; }
.card-primary { border-color: #cbd5e1 !important; border-top: 4px solid #2563eb !important; }
.card-neutral { border-color: #e2e8f0 !important; border-top: 4px solid #94a3b8 !important; }
.card-success { background: #f0fdf4 !important; border-color: #bbf7d0 !important; border-top: 4px solid #16a34a !important; }
.card-danger { background: #fef2f2 !important; border-color: #fecaca !important; border-top: 4px solid #dc2626 !important; }
.card-small { min-height: 90px !important; padding: 0.8rem 0.5rem !important; border-radius: 10px !important; }
.card-small .kpi-label { font-size: 0.72rem !important; margin-bottom: 0.2rem !important; }
.card-small .kpi-value { font-size: 1.45rem !important; }

.login-wrap { max-width:460px; margin:5rem auto 0; padding:2.8rem 2.5rem 1.5rem; background: #ffffff; border:1px solid #cbd5e1; border-radius:16px; box-shadow: 0 10px 15px -3px rgba(0,0,0,0.1), 0 4px 6px -2px rgba(0,0,0,0.05); }
.login-title { font-size: 2rem !important; font-weight: 800 !important; text-align: center; margin-bottom: .4rem; color: #1e293b; }
.login-sub { font-size: 1rem !important; text-align: center; color: #64748b; margin-bottom: 2rem; font-weight: 600 !important; }

.badge { display:inline-block; font-size:.72rem !important; border-radius:6px; padding:3px 11px; margin-left:6px; font-weight:800; letter-spacing:.05em; }
.badge-admin  { background:#eff6ff; border:1px solid #bfdbfe; color:#1d4ed8; }
.badge-expert { background:#f0fdf4; border:1px solid #bbf7d0; color:#166534; }
.badge-supervisor { background:#fefce8; border:1px solid #fef08a; color:#a16207; }
.req-pending  { background:#fffbeb; border:1px solid #fde047; border-radius:8px; padding:.85rem 1.4rem; margin-bottom:.8rem; color:#92400e; font-size:.95rem; font-weight:600; }
.section-card { background:#ffffff; border:1px solid #e2e8f0; border-radius:12px; padding:1.6rem 2rem; margin-bottom:1.4rem; box-shadow: 0 1px 2px rgba(0,0,0,0.05); }

.scorecard-container { width: 100%; overflow-x: auto; margin-top: 1rem; margin-bottom: 2rem; border-radius: 8px; box-shadow: 0 4px 6px rgba(0,0,0,0.05); border: 1px solid #e2e8f0; }
.scorecard-container table { width: 100%; border-collapse: collapse; background-color: #ffffff; font-family: inherit; table-layout: auto; }
.scorecard-container th { background-color: #f8fafc !important; color: #334155 !important; font-weight: 800 !important; font-size: 0.9rem !important; text-align: center !important; padding: 10px 8px !important; border-bottom: 2px solid #cbd5e1; }
.scorecard-container td { text-align: center !important; padding: 10px 8px !important; font-size: 0.95rem !important; border-bottom: 1px solid #f1f5f9; color: #1e293b; }
.scorecard-container td:first-child { font-weight: 700 !important; text-align: left !important; padding-left: 12px !important; }

.profile-card { background: #ffffff; border: 1px solid #e2e8f0; border-radius: 16px; padding: 1.5rem; text-align: center; box-shadow: 0 4px 6px -1px rgba(0,0,0,0.05); margin-bottom: 1.5rem; transition: transform 0.2s; }
.profile-card:hover { transform: translateY(-5px); box-shadow: 0 10px 15px -3px rgba(0,0,0,0.1); border-color: #94a3b8; }
.profile-img { width: 110px; height: 110px; border-radius: 50%; object-fit: cover; object-position: top; border: 3px solid #2563eb; margin: 0 auto 1rem; padding: 2px; }
.profile-name { font-size: 1.3rem !important; font-weight: 800 !important; color: #0f172a; margin-bottom: 2px; }
.profile-role { font-size: 0.8rem !important; color: #64748b; margin-bottom: 1.2rem; text-transform: uppercase; letter-spacing: 1px; font-weight: 700 !important; }
.profile-detail { font-size: 0.95rem !important; color: #334155; margin-bottom: 0.4rem; text-align: left; display:flex; justify-content: space-between; border-bottom: 1px dashed #f1f5f9; padding-bottom:4px; font-weight: 600 !important;}
.profile-bio { font-size: 0.95rem !important; font-style: italic; color: #475569; margin-top: 1rem; padding-top: 1rem; border-top: 1px solid #e2e8f0; font-weight: 500 !important;}
</style>
""", unsafe_allow_html=True)

THEME = dict(template="plotly_white", paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", font_color="#1e293b", margin=dict(l=10, r=10, t=55, b=10))

# ══════════════════════════════════════════════════════════════════════════════════
#  SESSION STATE & SECURE PERSISTENT LOGIN (WITH EXPIRY)
# ══════════════════════════════════════════════════════════════════════════════════
SESSION_DURATION_HOURS = 12  

def generate_signed_token(uname: str, exp_timestamp: str) -> str:
    return hashlib.sha256(f"{uname}|{exp_timestamp}|SECRET_ALDAWAA_TOKEN".encode()).hexdigest()

if "store" not in st.session_state: st.session_state.store = _load_store()
    
if "authenticated" not in st.session_state:
    st.session_state.authenticated, st.session_state.username, st.session_state.role = False, None, None
    if "usr" in st.query_params and "tok" in st.query_params and "exp" in st.query_params:
        q_usr, q_tok, q_exp = st.query_params["usr"], st.query_params["tok"], st.query_params["exp"]
        udata = st.session_state.store["users"].get(q_usr)
        try: is_expired = int(time.time()) > int(q_exp)
        except ValueError: is_expired = True
        if udata and not is_expired and generate_signed_token(q_usr, q_exp) == q_tok:
            st.session_state.authenticated, st.session_state.username, st.session_state.role = True, q_usr, udata["role"]
        else: st.query_params.clear()
            
if "page" not in st.session_state: st.session_state.page = "dashboard"
if "force_onboard" not in st.session_state: st.session_state.force_onboard = False
if "view_request_form" not in st.session_state: st.session_state.view_request_form = False

def users()     -> dict: return st.session_state.store["users"]
def requests()  -> list: return st.session_state.store["requests"]
def overrides() -> dict: return st.session_state.store["overrides"]
def me()        -> str:  return st.session_state.username
def is_admin()  -> bool: return st.session_state.role == "admin"
def cur_user()  -> dict: return users().get(me(), {})
def agent_name_of(uname: str) -> str: return users().get(uname, {}).get("agent_name")
def my_agent_name() -> str: return agent_name_of(me())
def pending_count() -> int: return sum(1 for r in requests() if r["status"] == "pending")

def push_request(uname, rtype, new_value):
    requests().append({"id": int(time.time() * 1000), "requester": uname, "type": rtype, "new_value": new_value, "status": "pending", "ts": time.strftime("%Y-%m-%d %H:%M")})
    _save_store()

def approve_request(req_id):
    for r in requests():
        if r["id"] == req_id and r["status"] == "pending":
            u = r["requester"]
            if r["type"] == "display_name": users()[u]["display_name"] = r["new_value"]
            elif r["type"] == "password":   
                users()[u]["password_hash"] = _hash(r["new_value"]); update_sheet_password(u, r["new_value"]) 
            r["status"] = "approved"; _save_store(); return True
    return False

def reject_request(req_id):
    for r in requests():
        if r["id"] == req_id and r["status"] == "pending": r["status"] = "rejected"; _save_store(); return True
    return False

def calc_change(curr, prev):
    if pd.isna(curr): curr = 0
    if pd.isna(prev): prev = 0
    if prev == 0: return 100.0 if curr > 0 else 0.0
    return ((curr - prev) / prev) * 100.0

def kpi_colored(label, value, cls, change=None, inverse=False, neutral=False):
    change_html = ""
    if change is not None:
        if neutral: color, arrow = "#64748b", ("▲" if change > 0 else ("▼" if change < 0 else "−"))
        else:
            if change > 0: arrow, color = "▲", ("#dc2626" if inverse else "#16a34a")
            elif change < 0: arrow, color = "▼", ("#16a34a" if inverse else "#dc2626")
            else: arrow, color = "−", "#94a3b8"
        bg_color = color + "15"
        change_html = f'<div style="position: absolute; bottom: 8px; right: 10px; font-size: 0.8rem; padding: 2px 8px; border-radius: 12px; font-weight: 700; color: {color}; background-color: {bg_color};">{arrow} {abs(change):.1f}%</div>'
    return (f'<div class="kpi-container {cls}"><div class="kpi-label">{label}</div><div class="kpi-value" style="display:flex; align-items:center; justify-content:center; flex-wrap:wrap; gap:6px;">{value}</div>{change_html}</div>')

def time_to_minutes(s):
    try: p = str(s).strip().split(":"); return int(p[0]) * 60 + int(p[1])
    except: return 0

def fmt_m(v):
    if pd.isna(v) or v <= 0: return "00:00:00"
    t = int(round(v * 60))
    return f"{t // 3600:02d}:{(t % 3600) // 60:02d}:{t % 60:02d}"

def assign_time_tier(m):
    if m <= 15: return "Under 15 Mins"
    if m <= 30: return "15-30 Mins"
    if m <= 45: return "30-45 Mins"
    if m <= 60: return "45-60 Mins"
    return "Over 1 Hour"

DAYS_AR = {"Saturday": "السبت", "Sunday": "الأحد", "Monday": "الإثنين", "Tuesday": "الثلاثاء", "Wednesday": "الأربعاء", "Thursday": "الخميس", "Friday": "الجمعة"}

OFFICIAL_EXPERTS = ["Ahmed El-Kholy", "Ahmed Kadry", "Amr El-Sayed", "Eslam Ramadan", "Mohamed Abdelmageed", "Mohamed Khalifa", "Yahia Ali Shafei"]
EXPERT_ID_MAP = {"Ahmed El-Kholy": "50107", "Ahmed Kadry": "50399", "Amr El-Sayed": "50187", "Eslam Ramadan": "50461", "Mohamed Abdelmageed": "50274", "Mohamed Khalifa": "50476", "Yahia Ali Shafei": "50114"}
AGENT_ALIASES = {"mohamed abdelmajid": "Mohamed Abdelmageed", "mohamed el-sayed": "Mohamed Abdelmageed", "محمد عبد المجيد": "Mohamed Abdelmageed", "محمد السيد عبد المجيد": "Mohamed Abdelmageed", "محمد السيد": "Mohamed Abdelmageed", "50274": "Mohamed Abdelmageed", "احمد الخولى": "Ahmed El-Kholy", "أحمد الخولي": "Ahmed El-Kholy", "احمد الخولي": "Ahmed El-Kholy", "50107": "Ahmed El-Kholy", "يحي علي شافعي": "Yahia Ali Shafei", "يحيي علي شافعي": "Yahia Ali Shafei", "50114": "Yahia Ali Shafei", "عمرو محمد السيد": "Amr El-Sayed", "50187": "Amr El-Sayed", "أحمد محمد قدري": "Ahmed Kadry", "احمد محمد قدري": "Ahmed Kadry", "احمد قدري": "Ahmed Kadry", "50399": "Ahmed Kadry", "إسلام رمضان خليل": "Eslam Ramadan", "أصلان رمضان خليل": "Eslam Ramadan", "اسلام رمضان": "Eslam Ramadan", "50461": "Eslam Ramadan", "محمد خليفة جاب الله": "Mohamed Khalifa", "محمد خليفه جاب الله": "Mohamed Khalifa", "محمد خليفة": "Mohamed Khalifa", "محمد خليفه": "Mohamed Khalifa", "50476": "Mohamed Khalifa", "محمد شحته عبدالنبي مصطفي": "Muhammad Shehta", "50228": "Muhammad Shehta"}
EXCLUSION_LIST = ['off', 'اوف', 'أوف', 'راحة', 'annual', 'casual', 'عارضة', 'عارضه', 'v', 'a', 'vacation', 'resign', 'استقالة', 'مستقيل', 'sick', 'مرضي', 'nan', 'none', '']

def normalize_expert_name(name):
    if pd.notna(name): name = re.sub(r'^\d+\s*-\s*', '', str(name).strip())
    n_lower = str(name).lower().strip()
    if n_lower in AGENT_ALIASES: return AGENT_ALIASES[n_lower]
    id_to_name = {v: k for k, v in EXPERT_ID_MAP.items()}
    if name in id_to_name: return id_to_name[name]
    return str(name).strip()

# ══════════════════════════════════════════════════════════════════════════════════
#  LOGIN GATE & FIRST-TIME LOGIN ONBOARDING
# ══════════════════════════════════════════════════════════════════════════════════
if not st.session_state.authenticated:
    st.markdown("""<div class='login-wrap'><div class='login-title'>💊 Dashboard Login</div><div class='login-sub'>In-Store Requests · AlDawaa</div></div>""", unsafe_allow_html=True)
    _, lc, _ = st.columns([1, 1.4, 1])
    with lc:
        if not st.session_state.view_request_form:
            inp_u = st.text_input("Username / ID", placeholder="Enter ID", key="li_u")
            inp_p = st.text_input("Password", type="password", placeholder="Enter password", key="li_p")
            if st.button("🔐 Login", use_container_width=True):
                uname = inp_u.strip().lower()
                udata = users().get(uname)
                if udata and udata["password_hash"] == _hash(inp_p):
                    notify_admin_whatsapp(udata.get("display_name", uname) + " ✅ Success")
                    if "login_logs" not in st.session_state.store: st.session_state.store["login_logs"] = []
                    st.session_state.store["login_logs"].append({"Timestamp": time.strftime('%Y-%m-%d %H:%M:%S'), "Username": uname, "Display Name": udata.get("display_name", uname), "Role": udata["role"]})
                    st.session_state.store["login_logs"] = st.session_state.store["login_logs"][-200:]
                    _save_store()
                    exp_time = str(int(time.time()) + (SESSION_DURATION_HOURS * 3600))
                    st.query_params["usr"] = uname
                    st.query_params["exp"] = exp_time
                    st.query_params["tok"] = generate_signed_token(uname, exp_time)
                    if inp_u.strip() == inp_p.strip() and udata["role"] != "admin":
                        st.session_state.username = uname
                        st.session_state.force_onboard = True
                        st.session_state.authenticated = True
                        st.rerun()
                    else:
                        st.session_state.authenticated = True
                        st.session_state.username = uname
                        st.session_state.role = udata["role"]
                        st.session_state.page = "dashboard"
                        st.rerun()
                else:
                    notify_admin_whatsapp(uname + " ❌ Failed Attempt")
                    st.error("❌ Incorrect username or password.")
            st.write("")
            if st.button("🆕 Create Account / Request Access", use_container_width=True):
                st.session_state.view_request_form = True; st.rerun()
        else:
            st.markdown("### 📝 Request Account Creation")
            req_name = st.text_input("Full Name *", placeholder="e.g. Ahmed Ali")
            req_id = st.text_input("Username / ID *", placeholder="e.g. 50123")
            req_email = st.text_input("Email *", placeholder="e.g. ahmed@example.com") 
            if st.button("📤 Submit Access Request", use_container_width=True):
                if req_name.strip() and req_id.strip() and req_email.strip():
                    uid = req_id.strip().lower()
                    if uid in users(): st.error("❌ This Username/ID is already registered.")
                    else:
                        payload = json.dumps({"name": req_name.strip(), "id": uid, "email": req_email.strip()})
                        push_request(uid, "new_account", payload)
                        st.success(f"✅ Request sent! Waiting for admin approval. You will receive an email upon approval."); time.sleep(2.5)
                        st.session_state.view_request_form = False; st.rerun()
                else: st.error("❌ Name, ID, and Email fields are all required.")
            if st.button("← Back to Login", use_container_width=True): st.session_state.view_request_form = False; st.rerun()
    st.stop()

if st.session_state.force_onboard:
    st.markdown("## ⚙️ Mandatory Password Update Required")
    st.info("🚨 This is your first login. You must update your password before accessing dashboard metrics.")
    _, ob_col, _ = st.columns([1, 1.5, 1])
    with ob_col:
        st.markdown("<div class='section-card'>", unsafe_allow_html=True)
        with st.form("onboard_pass_form"):
            new_ob1 = st.text_input("Set New Password", type="password")
            new_ob2 = st.text_input("Confirm New Password", type="password")
            if st.form_submit_button("💾 Save & Open Dashboard", use_container_width=True):
                if new_ob1 != new_ob2: st.error("❌ Passwords do not match.")
                elif len(new_ob1) < 6: st.error("❌ Password must be at least 6 characters.")
                else:
                    uname = st.session_state.username
                    users()[uname]["password_hash"] = _hash(new_ob1)
                    _save_store(); update_sheet_password(uname, new_ob1) 
                    st.session_state.role = users()[uname]["role"]
                    st.session_state.force_onboard = False
                    st.session_state.page = "dashboard"
                    st.success("✅ Password configured successfully!"); time.sleep(1.5); st.rerun()
        st.markdown("</div>", unsafe_allow_html=True)
    st.stop()

# ══════════════════════════════════════════════════════════════════════════════════
#  SIDEBAR MODULE
# ══════════════════════════════════════════════════════════════════════════════════
with st.sidebar:
    st.markdown("## Approvals Team Dashboard")
    if is_admin(): badge_cls, badge_txt = "badge-admin", "ADMIN"
    elif st.session_state.role == "supervisor": badge_cls, badge_txt = "badge-supervisor", "SUPERVISOR"
    else: badge_cls, badge_txt = "badge-expert", "EXPERT"
    st.markdown(f"👤 **{cur_user().get('display_name', '–')}** <span class='badge {badge_cls}'>{badge_txt}</span>", unsafe_allow_html=True)

    @st.cache_data(ttl=600, show_spinner="Syncing database tables…")
    def load_data():
        try:
            scopes = ["https://www.googleapis.com/auth/spreadsheets", "https://www.googleapis.com/auth/drive"]
            if "gspread" in st.secrets and "credentials" in st.secrets["gspread"]:
                creds = Credentials.from_service_account_info(json.loads(st.secrets["gspread"]["credentials"]), scopes=scopes)
            else:
                st.error("❌ Secrets file layout unconfigured.")
                return pd.DataFrame(), pd.DataFrame(), pd.DataFrame(), pd.DataFrame(), pd.DataFrame()
            
            client = gspread.authorize(creds)
            sheet = client.open("AlDawaa Tickets Data")
            all_dfs, roster_df, out_req_df, df_quality, df_users = [], pd.DataFrame(), pd.DataFrame(), pd.DataFrame(), pd.DataFrame()
            for ws in sheet.worksheets():
                title = ws.title.strip()
                data = ws.get_all_values()
                if len(data) < 2: continue
                if title == "Working Days": roster_df = pd.DataFrame(data[1:], columns=[str(c).strip() for c in data[0]]); continue
                elif title == "Out Requests": out_req_df = pd.DataFrame(data[1:], columns=[str(c).strip() for c in data[0]]); continue
                elif title == "Quality Issues": df_quality = pd.DataFrame(data[1:], columns=[str(c).strip() for c in data[0]]); continue
                elif title == "Users": 
                    df_users = pd.DataFrame(data[1:], columns=[str(c).strip() for c in data[0]])
                    continue
                
                dft = pd.DataFrame(data[1:], columns=[str(c).strip() for c in data[0]])
                mp = {}; seen = set()
                for col in dft.columns:
                    cl = col.lower(); t = None
                    if "id" in cl and "req" in cl: t = "Request ID"
                    elif "date" in cl: t = "Request Date"
                    elif "type" in cl: t = "Request Type"
                    elif "status" in cl and "count" in cl: t = "Status Count"
                    elif "status" in cl: t = "Status"
                    elif "assigned" in cl or "agent" in cl: t = "Assigned By"
                    elif "response" in cl and "take" in cl: t = "Response Take"
                    elif "action" in cl and "take" in cl: t = "First Action Take"
                    elif "request" in cl and "take" in cl: t = "Request Take"
                    elif "email" in cl or "special" in cl: t = "Is Special Request(By Email)"
                    elif "hic" in cl or "insurance" in cl: t = "HIC"
                    elif "store" in cl or "branch" in cl or "pharmacy" in cl: t = "Store ID"
                    if t and t not in seen: mp[col] = t; seen.add(t)
                dft.rename(columns=mp, inplace=True)
                all_dfs.append(dft)
            
            if not all_dfs: return pd.DataFrame(), roster_df, out_req_df, df_quality, df_users
                
            df = pd.concat(all_dfs, ignore_index=True, sort=False).replace("", np.nan)
            for c in ["Request ID", "Request Date", "Request Type", "Status", "Status Count", "Request Take", "Response Take", "First Action Take", "Assigned By", "Is Special Request(By Email)", "HIC"]:
                if c not in df.columns: df[c] = np.nan
            if "Store ID" not in df.columns: df["Store ID"] = "Unknown"
            
            df["Status"]       = df["Status"].fillna("Unknown")
            df["Status Count"] = pd.to_numeric(df["Status Count"], errors="coerce").fillna(0).astype(int)
            df["Request Type"] = df["Request Type"].fillna("Unknown Type")
            df["HIC"]          = df["HIC"].fillna("Unknown")
            df["Assigned By"]  = df["Assigned By"].fillna("Unassigned").astype(str).str.strip().apply(normalize_expert_name)
            df["Store ID"]     = df["Store ID"].fillna("Unknown").astype(str).str.strip()
            
            dp = pd.to_datetime(df["Request Date"], errors="coerce")
            df["Request Date"], df["Date Only"], df["Hour"], df["Day Name"] = dp, dp.dt.date, dp.dt.hour.fillna(0).astype(int), dp.dt.day_name().fillna("Unknown")
            if "Request Take" in df.columns: df["Request Take (min)"] = df["Request Take"].apply(time_to_minutes).fillna(0)
            if "Response Take" in df.columns: df["Response Take (min)"] = df["Response Take"].apply(time_to_minutes).fillna(0)
            df["Is Email"] = (df["Is Special Request(By Email)"].astype(str).str.strip().str.lower() == "yes")
                
            return df, roster_df, out_req_df, df_quality, df_users
        except Exception as e:
            st.error(f"❌ Connection Error: {e}")
            return pd.DataFrame(), pd.DataFrame(), pd.DataFrame(), pd.DataFrame(), pd.DataFrame()

    df_raw, df_roster, df_out_req, df_quality, df_users = load_data()
    
    # ── SYNC USERS PROFILES FROM SHEET TO DASHBOARD ──
    if not df_users.empty:
        updated = False
        for _, row in df_users.iterrows():
            if len(row) > 0 and pd.notna(row.iloc[0]):
                uname = str(row.iloc[0]).strip().lower()
                if uname in st.session_state.store["users"]:
                    u_dict = st.session_state.store["users"][uname]
                    if len(row) > 6 and str(row.iloc[6]).strip(): u_dict["photo"] = parse_drive_link(str(row.iloc[6]))
                    if len(row) > 7 and str(row.iloc[7]).strip(): u_dict["grad_year"] = str(row.iloc[7]).strip()
                    if len(row) > 8 and str(row.iloc[8]).strip(): u_dict["join_cc"] = str(row.iloc[8]).strip()
                    if len(row) > 9 and str(row.iloc[9]).strip(): u_dict["join_team"] = str(row.iloc[9]).strip()
                    if len(row) > 10 and str(row.iloc[10]).strip(): u_dict["bio"] = str(row.iloc[10]).strip()
                    updated = True
        if updated: _save_store()

    if df_raw.empty: st.warning("Empty source records."); st.stop()

    raw_dates = pd.to_datetime(df_raw["Request Date"]).dropna()
    if not raw_dates.empty:
        max_uploaded_date = raw_dates.max()
        first_of_current_upload_month = max_uploaded_date.replace(day=1)
        last_day_of_ended_month = first_of_current_upload_month - pd.Timedelta(days=1)
        default_from, default_to = last_day_of_ended_month.replace(day=1).date(), last_day_of_ended_month.date()
    else:
        default_from, default_to = df_raw["Date Only"].dropna().min(), df_raw["Date Only"].dropna().max()

    min_d, max_d = df_raw["Date Only"].dropna().min(), df_raw["Date Only"].dropna().max()
    date_range = st.date_input("Date Range", value=(default_from, default_to), min_value=min_d, max_value=max_d)
    d_from, d_to = date_range if isinstance(date_range, (list, tuple)) and len(date_range) == 2 else (min_d, max_d)
    if d_from == d_to: st.caption(f"📅 {DAYS_AR.get(pd.to_datetime(d_from).day_name(), '')}")
        
    delta_days = (d_to - d_from).days + 1
    prev_d_from, prev_d_to = d_from - timedelta(days=delta_days), d_from - timedelta(days=1)
    PERIOD_KEY = f"{d_from}_{d_to}"
    
    sel_hic = st.multiselect("HIC", sorted(df_raw["HIC"].dropna().unique()))
    sel_req_type = st.multiselect("Request Type", sorted(df_raw["Request Type"].dropna().unique()))

    st.markdown("<br><br>", unsafe_allow_html=True); st.divider()
    if is_admin() and pending_count() > 0: st.warning(f"🔔 {pending_count()} pending requests")
    st.success("📡 Live Sync Active")
    if is_admin() and st.button("🔄 Refresh Data Now", use_container_width=True): st.cache_data.clear(); st.rerun()

    sb1, sb2 = st.columns(2)
    with sb1:
        if st.button("⚙️ Settings", use_container_width=True): st.session_state.page = "settings"; st.rerun()
    with sb2:
        if st.button("🚪 Logout", use_container_width=True):
            st.query_params.clear()
            for k in ("authenticated", "username", "role", "page", "force_onboard"): st.session_state.pop(k, None)
            st.rerun()

# ══════════════════════════════════════════════════════════════════════════════════
#  HELPER FUNCTIONS (ROSTER & PARSING)
# ══════════════════════════════════════════════════════════════════════════════════
def get_roster_stats(roster_df, start_date, end_date):
    counts = {exp: {"wd": 0, "off": 0, "ann": 0, "cas": 0, "sick": 0} for exp in OFFICIAL_EXPERTS}
    if roster_df.empty: return counts
    valid_cols = []
    for col in roster_df.columns:
        col_str = str(col).strip()
        clean_col = re.sub(r'(?i)\b(monday|tuesday|wednesday|thursday|friday|saturday|sunday)\b', '', col_str).strip()
        match = re.search(r'(\d{1,2}[-/\s]+(?:[A-Za-z]+|\d{1,2})[-/\s]+\d{4}|\d{4}[-/\s]+\d{1,2}[-/\s]+\d{1,2})', clean_col)
        if match:
            try:
                col_date = pd.to_datetime(match.group(1), dayfirst=True).date()
                if start_date <= col_date <= end_date: valid_cols.append(col)
            except: pass

    df_r_str = roster_df.astype(str)
    for exp in OFFICIAL_EXPERTS:
        exp_id = EXPERT_ID_MAP.get(exp, "")
        if exp_id:
            mask = df_r_str.apply(lambda row: str(exp_id).strip().lower() in [str(x).strip().lower() for x in row.values], axis=1)
            exp_rows = df_r_str[mask]
            if not exp_rows.empty and valid_cols:
                safe_cols = [c for c in valid_cols if c in exp_rows.columns]
                if safe_cols:
                    vals_lower = [str(v).strip().lower() for v in exp_rows[safe_cols].values.flatten()]
                    counts[exp]["off"] = sum(1 for v in vals_lower if v in ['off', 'اوف', 'أوف', 'راحة'])
                    counts[exp]["ann"] = sum(1 for v in vals_lower if v in ['annual', 'v', 'a', 'vacation'])
                    counts[exp]["cas"] = sum(1 for v in vals_lower if v in ['casual', 'عارضة', 'عارضه'])
                    counts[exp]["sick"] = sum(1 for v in vals_lower if v in ['sick', 'مرضي', 'مرضى'])
                    counts[exp]["wd"] = sum(1 for v in vals_lower if v and v not in EXCLUSION_LIST)
    return counts

curr_roster_counts = get_roster_stats(df_roster, d_from, d_to)
prev_roster_counts = get_roster_stats(df_roster, prev_d_from, prev_d_to)

out_req_dict, prev_out_req_dict = {}, {}
global_jhah, global_support, prev_global_jhah, prev_global_support = 0, 0, 0, 0

if not df_out_req.empty and "Date" in df_out_req.columns:
    df_out_req["Parsed_Date"] = pd.to_datetime(df_out_req["Date"], errors="coerce").dt.date
    
    df_out_filtered = df_out_req[(df_out_req["Parsed_Date"] >= d_from) & (df_out_req["Parsed_Date"] <= d_to)].copy()
    if "Source" in df_out_filtered.columns: global_jhah = df_out_filtered[df_out_filtered["Source"].astype(str).str.lower().str.contains("jhah", na=False)].shape[0]
    global_support = len(df_out_filtered) - global_jhah
    if "Expert Name" in df_out_filtered.columns:
        df_out_filtered["Norm_Expert"] = df_out_filtered["Expert Name"].fillna("").astype(str).apply(normalize_expert_name).str.lower()
        for exp_name, grp in df_out_filtered.groupby("Norm_Expert"):
            if not exp_name or exp_name == "nan": continue 
            j_count = grp[grp["Source"].astype(str).str.lower().str.contains("jhah", na=False)].shape[0] if "Source" in grp.columns else 0
            out_req_dict[exp_name] = {"JHAH": j_count, "Support Req": len(grp) - j_count}

    df_out_prev = df_out_req[(df_out_req["Parsed_Date"] >= prev_d_from) & (df_out_req["Parsed_Date"] <= prev_d_to)].copy()
    if "Source" in df_out_prev.columns: prev_global_jhah = df_out_prev[df_out_prev["Source"].astype(str).str.lower().str.contains("jhah", na=False)].shape[0]
    prev_global_support = len(df_out_prev) - prev_global_jhah
    if "Expert Name" in df_out_prev.columns:
        df_out_prev["Norm_Expert"] = df_out_prev["Expert Name"].fillna("").astype(str).apply(normalize_expert_name).str.lower()
        for exp_name, grp in df_out_prev.groupby("Norm_Expert"):
            if not exp_name or exp_name == "nan": continue 
            j_count = grp[grp["Source"].astype(str).str.lower().str.contains("jhah", na=False)].shape[0] if "Source" in grp.columns else 0
            prev_out_req_dict[exp_name] = {"JHAH": j_count, "Support Req": len(grp) - j_count}

df = df_raw[(df_raw["Date Only"] >= d_from) & (df_raw["Date Only"] <= d_to)].copy()
df_prev_all = df_raw[(df_raw["Date Only"] >= prev_d_from) & (df_raw["Date Only"] <= prev_d_to)].copy()

if sel_hic: df = df[df["HIC"].isin(sel_hic)]; df_prev_all = df_prev_all[df_prev_all["HIC"].isin(sel_hic)]
if sel_req_type: df = df[df["Request Type"].isin(sel_req_type)]; df_prev_all = df_prev_all[df_prev_all["Request Type"].isin(sel_req_type)]

expert_quality_deductions, df_q_filtered = {}, pd.DataFrame()
if not df_quality.empty and "Date" in df_quality.columns and "Severity" in df_quality.columns and "Expert Name" in df_quality.columns:
    df_quality['Parsed_Date'] = pd.to_datetime(df_quality['Date'], errors='coerce').dt.date
    df_q_filtered = df_quality[(df_quality['Parsed_Date'] >= d_from) & (df_quality['Parsed_Date'] <= d_to)].copy()
    def get_deduction(sev):
        s = str(sev).strip().lower()
        if s == 'critical': return 5.0
        elif s == 'major': return 2.0
        elif s == 'medium': return 1.0
        elif s == 'minor': return 0.5
        return 0.0
    df_q_filtered['Deduction'] = df_q_filtered['Severity'].apply(get_deduction)
    df_q_filtered['Norm_Expert'] = df_q_filtered['Expert Name'].apply(normalize_expert_name).str.lower()
    df_q_filtered['Display_Expert'] = df_q_filtered['Expert Name'].apply(normalize_expert_name)
    expert_quality_deductions = df_q_filtered.groupby('Norm_Expert')['Deduction'].sum().to_dict()

# ══════════════════════════════════════════════════════════════════════════════════
#  SETTINGS PANEL (INCLUDES PROFILE EDIT FORM)
# ══════════════════════════════════════════════════════════════════════════════════
if st.session_state.page == "settings":
    if st.button("← Back to Dashboard"): st.session_state.page = "dashboard"; st.rerun()

    def render_profile_edit_form(urow):
        st.markdown("<div class='section-card'>", unsafe_allow_html=True)
        st.markdown("### 📋 Edit My Professional Profile")
        with st.form("expert_profile_form"):
            st.markdown("**📸 Profile Picture**")
            
            curr_photo = urow.get("photo", "")
            disp_url = curr_photo if not curr_photo.startswith("data:image") else ""
            
            p_photo_file = st.file_uploader("1️⃣ Upload from Phone/Computer", type=["png", "jpg", "jpeg"])
            p_photo_url = st.text_input("2️⃣ OR Paste Google Drive/Image Link", value=disp_url)
            
            c1, c2 = st.columns(2)
            with c1: p_grad = st.text_input("🎓 Graduation Year", value=urow.get("grad_year", ""))
            with c2: p_join_cc = st.text_input("🏢 Joined Call Center (Year/Month)", value=urow.get("join_cc", ""))
            c3, c4 = st.columns(2)
            with c3: p_join_tm = st.text_input("🚀 Joined Approvals Team (Year/Month)", value=urow.get("join_team", ""))
            with c4: p_bio = st.text_input("✍️ Short Bio / Quote", value=urow.get("bio", ""))
            
            if st.form_submit_button("💾 Save Profile Data", use_container_width=True):
                final_photo = curr_photo
                if p_photo_file is not None:
                    try:
                        img = Image.open(p_photo_file).convert("RGB")
                        img.thumbnail((250, 250)) 
                        buffered = BytesIO()
                        img.save(buffered, format="JPEG", quality=80)
                        img_str = base64.b64encode(buffered.getvalue()).decode()
                        final_photo = f"data:image/jpeg;base64,{img_str}"
                    except: pass
                elif p_photo_url.strip() and p_photo_url.strip() != disp_url:
                    final_photo = parse_drive_link(p_photo_url.strip())
                
                users()[me()]["photo"] = final_photo
                users()[me()]["grad_year"] = p_grad.strip()
                users()[me()]["join_cc"] = p_join_cc.strip()
                users()[me()]["join_team"] = p_join_tm.strip()
                users()[me()]["bio"] = p_bio.strip()
                _save_store()
                update_sheet_profile(me(), users()[me()])
                st.success("✅ Profile Updated successfully!"); time.sleep(1); st.rerun()
        st.markdown("</div>", unsafe_allow_html=True)

    if is_admin():
        st.markdown("## ⚙️ Admin Control Panel")
        atab1, atab2, atab3, atab4 = st.tabs(["👤 My Profile", "🔔 Change & Access Requests", "👥 Manage Dashboard Users", "📜 System Access Logs"])

        with atab1:
            urow = users()[me()]
            st.markdown("<div class='section-card'>", unsafe_allow_html=True)
            st.markdown("**✏️ Display Name**")
            with st.form("admin_name_form"):
                new_dn = st.text_input("New Display Name", value=urow["display_name"])
                if st.form_submit_button("💾 Save Name", use_container_width=True):
                    if new_dn.strip() and new_dn.strip() != urow["display_name"]:
                        users()[me()]["display_name"] = new_dn.strip(); _save_store(); st.success("✅ Profile display name updated."); st.rerun()
            st.markdown("</div>", unsafe_allow_html=True)
            
            render_profile_edit_form(urow)

            st.markdown("<div class='section-card'>", unsafe_allow_html=True)
            st.markdown("**🔑 Change Password**")
            with st.form("admin_pw_form"):
                old_pw  = st.text_input("Current Account Password", type="password")
                new_pw1 = st.text_input("New Secure Password",     type="password")
                new_pw2 = st.text_input("Confirm Secure Password", type="password")
                if st.form_submit_button("💾 Update Password", use_container_width=True):
                    if _hash(old_pw) != urow["password_hash"]: st.error("❌ Password input incorrect.")
                    elif new_pw1 != new_pw2: st.error("❌ Confirmation inputs mismatch.")
                    elif len(new_pw1) < 6:   st.error("❌ Minimum limit 6 characters.")
                    else:
                        users()[me()]["password_hash"] = _hash(new_pw1)
                        _save_store(); update_sheet_password(me(), new_pw1); st.success("✅ Administrative password saved.")
            st.markdown("</div>", unsafe_allow_html=True)

        with atab2:
            st.markdown("### 🔔 Change & Access Requests")
            pending = [r for r in requests() if r["status"] == "pending"]
            if not pending: st.info("✅ No requests pending approval.")
            else:
                for req in pending:
                    if req["type"] == "new_account":
                        try: p_data = json.loads(req["new_value"])
                        except: p_data = {"name": req["requester"], "id": req["requester"], "email": ""}
                        st.markdown(f"<div class='req-pending'>🕐 <b>{req['ts']}</b> &nbsp;|&nbsp; 🆕 <b>NEW ACCOUNT REQUEST</b> <br>Name: <b>{p_data.get('name')}</b> &nbsp;|&nbsp; ID: <b>{p_data.get('id')}</b> &nbsp;|&nbsp; Email: <b>{p_data.get('email', 'N/A')}</b></div>", unsafe_allow_html=True)
                        rc1, rc2, rc3 = st.columns([3, 1, 1])
                        with rc2:
                            if st.button("✅ Approve", key=f"apr_na_{req['id']}", use_container_width=True):
                                users()[p_data['id']] = {"display_name": p_data['name'], "password_hash": _hash(p_data['id']), "role": "expert", "agent_name": p_data['name']}
                                req["status"] = "approved"; _save_store()
                                try:
                                    scopes = ["https://www.googleapis.com/auth/spreadsheets", "https://www.googleapis.com/auth/drive"]
                                    if "gspread" in st.secrets and "credentials" in st.secrets["gspread"]:
                                        creds = Credentials.from_service_account_info(json.loads(st.secrets["gspread"]["credentials"]), scopes=scopes)
                                        gspread.authorize(creds).open("AlDawaa Tickets Data").worksheet("Users").append_row([p_data['id'], p_data['id'], "expert", p_data['name'], p_data['name'], p_data.get('email', '')])
                                except Exception as e: pass 
                                user_email = p_data.get('email', '').strip()
                                if user_email:
                                    if send_approval_email(user_email, p_data['name'], p_data['id']): st.success("✅ Account approved & Email sent successfully!")
                                    else: st.warning("✅ Account approved, but Email failed (Check your SMTP Settings).")
                                else: st.success("✅ Account approved! (No valid email was provided)")
                                time.sleep(2.5); st.rerun()
                        with rc3:
                            if st.button("❌ Deny", key=f"rej_na_{req['id']}", use_container_width=True): reject_request(req["id"]); st.warning("Request rejected."); time.sleep(1); st.rerun()
                    else:
                        udata_r, udisp = users().get(req["requester"], {}), users().get(req["requester"], {}).get("display_name", req["requester"])
                        st.markdown(f"<div class='req-pending'>🕐 <b>{req['ts']}</b> &nbsp;|&nbsp; 👤 <b>{udisp}</b> &nbsp;|&nbsp; Wants to adjust <b>{'Display Name' if req['type'] == 'display_name' else 'Password'}</b></div>", unsafe_allow_html=True)
                        rc1, rc2, rc3 = st.columns([3, 1, 1])
                        with rc2:
                            if st.button("✅ Approve", key=f"apr_{req['id']}", use_container_width=True): approve_request(req["id"]); st.success("Approved successfully."); st.rerun()
                        with rc3:
                            if st.button("❌ Reject", key=f"rej_{req['id']}", use_container_width=True): reject_request(req["id"]); st.warning("Rejected successfully."); st.rerun()

        with atab3:
            st.markdown("### 👥 Manage Dashboard Users")
            with st.expander("➕ Add New User"):
                with st.form("add_new_user_form"):
                    c1, c2 = st.columns(2)
                    with c1:
                        add_uname = st.text_input("Username / ID *", key="add_uname")
                        add_dname = st.text_input("Display Name *", key="add_dname")
                        add_role = st.selectbox("Role", ["expert", "supervisor", "admin"], key="add_role")
                    with c2:
                        add_aname = st.text_input("Agent Key Mapping (Sheets)", key="add_aname")
                        add_pass1 = st.text_input("Password *", type="password", key="add_pass1")
                        add_pass2 = st.text_input("Confirm Password *", type="password", key="add_pass2")
                    if st.form_submit_button("➕ Create Account", use_container_width=True):
                        add_uname = add_uname.strip().lower()
                        if not add_uname or not add_dname.strip() or not add_pass1: st.error("❌ Please fill in all required fields (*).")
                        elif add_uname in users(): st.error("❌ Username/ID already exists.")
                        elif add_pass1 != add_pass2: st.error("❌ Passwords do not match.")
                        elif len(add_pass1) < 6: st.error("❌ Password must be at least 6 characters.")
                        else:
                            users()[add_uname] = {"display_name": add_dname.strip(), "password_hash": _hash(add_pass1), "role": add_role, "agent_name": add_aname.strip() if add_aname.strip() else None}
                            _save_store(); st.success(f"✅ Account for {add_dname.strip()} created successfully!"); time.sleep(1); st.rerun()
            st.divider()
            for uname, urow in list(users().items()):
                role_icon = "🔑" if urow["role"] == "admin" else ("👁️" if urow["role"] == "supervisor" else "👤")
                with st.expander(f"{role_icon} {urow['display_name']} (@{uname})"):
                    with st.form(f"admin_edit_{uname}"):
                        eu_dn   = st.text_input("Display Username", value=urow["display_name"], key=f"dn_{uname}")
                        eu_an   = st.text_input("Agent Key Mapping (Sheets)", value=urow.get("agent_name") or "", key=f"an_{uname}")
                        eu_p1   = st.text_input("Override Password", type="password", key=f"p1_{uname}")
                        eu_p2   = st.text_input("Confirm Password", type="password", key=f"p2_{uname}")
                        try: current_role_idx = ["expert", "supervisor", "admin"].index(urow["role"])
                        except: current_role_idx = 0
                        eu_role = st.selectbox("Role", ["expert", "supervisor", "admin"], index=current_role_idx, key=f"rl_{uname}")
                        col1, col2 = st.columns([3, 1])
                        with col1: saved = st.form_submit_button("💾 Update User Settings", use_container_width=True)
                        with col2: deleted = st.form_submit_button("🗑️ Delete User", use_container_width=True)
                    if saved:
                        if eu_dn.strip(): users()[uname]["display_name"] = eu_dn.strip()
                        users()[uname]["agent_name"] = eu_an.strip() if eu_an.strip() else None
                        users()[uname]["role"] = eu_role
                        if eu_p1 and eu_p1 == eu_p2: 
                            users()[uname]["password_hash"] = _hash(eu_p1)
                            update_sheet_password(uname, eu_p1)
                        _save_store(); st.success("✅ User settings updated."); st.rerun()
                    if deleted:
                        if uname == "admin": st.error("❌ Cannot delete the primary admin account!")
                        elif uname == me(): st.error("❌ You cannot delete your own account while logged in!")
                        else: users().pop(uname); _save_store(); st.success(f"🗑️ Account for {uname} has been successfully revoked and deleted."); time.sleep(1); st.rerun()
                            
        with atab4:
            st.markdown("### 📜 System Access Logs")
            logs = st.session_state.store.get("login_logs", [])
            if not logs: st.info("No login logs available yet.")
            else:
                df_logs = pd.DataFrame(logs).sort_values(by="Timestamp", ascending=False).reset_index(drop=True)
                st.dataframe(df_logs, use_container_width=True)
                if st.button("🗑️ Clear Logs", type="primary"): st.session_state.store["login_logs"] = []; _save_store(); st.success("Logs cleared successfully!"); st.rerun()

    else:
        st.markdown("## ⚙️ My Profile Settings")
        urow = users()[me()]
        st.markdown("<div class='section-card'>", unsafe_allow_html=True)
        st.markdown("### ✏️ Propose Profile Display Name Change")
        with st.form("expert_name_form"):
            req_name = st.text_input("Proposed Display Name", placeholder=urow["display_name"])
            if st.form_submit_button("📤 Submit Request", use_container_width=True):
                if req_name.strip(): push_request(me(), "display_name", req_name.strip()); st.success("✅ Request delivered."); st.rerun()
        st.markdown("</div>", unsafe_allow_html=True)
        
        render_profile_edit_form(urow)

        st.markdown("<div class='section-card'>", unsafe_allow_html=True)
        st.markdown("### 🔑 Direct Password Adjuster Form")
        with st.form("expert_pw_direct"):
            get_cur_pw = st.text_input("Verify Current Password", type="password")
            new_p1     = st.text_input("Set New Secret Password", type="password")
            new_p2     = st.text_input("Confirm New Secret Password", type="password")
            if st.form_submit_button("💾 Save Changes", use_container_width=True):
                if _hash(get_cur_pw) == urow["password_hash"] and new_p1 == new_p2 and len(new_p1) >= 6:
                    users()[me()]["password_hash"] = _hash(new_p1); _save_store(); update_sheet_password(me(), new_p1); st.success("✅ Password updated."); st.rerun()
        st.markdown("</div>", unsafe_allow_html=True)
    st.stop()

def generate_pptx_summary(d_from_str, d_to_str, total_val, avg_per_day_val, ok_val, ok_pct_val, issue_val, issue_pct_val, afr_str, tat_str, stores_val, actions_val, jhah_val, support_val, req_counts_dict, req_pct_dict, sc_g, fig_r, fig_st, fig_d, fig_hic):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        from io import BytesIO

        def add_plot_to_slide(prs, fig, title_str):
            if fig is None: return
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = title_str
            try:
                img_bytes = fig.to_image(format="png", width=950, height=500, engine="kaleido")
                image_stream = BytesIO(img_bytes)
                slide.shapes.add_picture(image_stream, Inches(0.5), Inches(1.5), width=Inches(9))
            except Exception as e:
                tf = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1)).text_frame
                tf.text = f"⚠️ Please add 'kaleido' to requirements.txt to render this chart.\n\nError details: {str(e)}"
                tf.paragraphs[0].font.color.rgb = RGBColor(200, 0, 0)

        prs = Presentation()

        # --- SLIDE 1: Title Slide ---
        slide_title = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide_title.shapes.title
        subtitle = slide_title.placeholders[1]
        title.text = "In-Store Requests: Executive Report"
        subtitle.text = f"Performance Period: {d_from_str} to {d_to_str}\nAlDawaa Approvals Team"

        # --- SLIDE 2: KPIs ---
        slide_kpi = prs.slides.add_slide(prs.slide_layouts[5])
        slide_kpi.shapes.title.text = "Operational Insights & KPIs"
        
        table_shape = slide_kpi.shapes.add_table(6, 4, Inches(0.5), Inches(1.5), Inches(9), Inches(3)).table
        table_shape.columns[0].width = Inches(2.5)
        table_shape.columns[1].width = Inches(2.0)
        table_shape.columns[2].width = Inches(2.5)
        table_shape.columns[3].width = Inches(2.0)
        metrics = [
            ("Total Requests", f"{total_val:,}"),
            ("Avg Requests / Day", f"{avg_per_day_val:.1f}"),
            ("Closed Completed", f"{ok_val:,} ({ok_pct_val:.1f}%)"),
            ("Closed with Issue", f"{issue_val:,} ({issue_pct_val:.1f}%)"),
            ("AFR (Avg Response)", afr_str),
            ("TAT (Avg Service)", tat_str),
            ("Stores Served", f"{stores_val:,}"),
            ("Total Actions", f"{actions_val:,}"),
            ("JHAH Requests", f"{jhah_val:,}"),
            ("Support Requests", f"{support_val:,}")
        ]
        for i in range(4):
            table_shape.cell(0, i).text = "Metric" if i % 2 == 0 else "Value"
            table_shape.cell(0, i).fill.solid()
            table_shape.cell(0, i).fill.fore_color.rgb = RGBColor(37, 99, 235)
            for paragraph in table_shape.cell(0, i).text_frame.paragraphs:
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.font.bold = True
                paragraph.alignment = PP_ALIGN.CENTER
        r, c = 1, 0
        for m_name, m_val in metrics:
            table_shape.cell(r, c).text = m_name
            table_shape.cell(r, c+1).text = str(m_val)
            for paragraph in table_shape.cell(r, c+1).text_frame.paragraphs:
                paragraph.font.bold = True
            c += 2
            if c >= 4: c, r = 0, r + 1

        # --- SLIDE 3: Team Scorecard ---
        slide_team = prs.slides.add_slide(prs.slide_layouts[5])
        slide_team.shapes.title.text = "Team Performance Scorecard"
        if not sc_g.empty:
            cols_to_show = ["Rank", "Expert", "Tickets Count", "Cases/Day", "% Achievement from Target", "Service Quality"]
            rows = len(sc_g) + 1
            cols = len(cols_to_show)
            table_team = slide_team.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.4 * rows)).table
            
            table_team.columns[0].width = Inches(0.8) # Rank
            table_team.columns[1].width = Inches(2.2) # Expert
            table_team.columns[2].width = Inches(1.5) # Tickets
            table_team.columns[3].width = Inches(1.5) # Cases/Day
            table_team.columns[4].width = Inches(1.5) # Target
            table_team.columns[5].width = Inches(1.5) # Quality
            
            for c_idx, col_name in enumerate(cols_to_show):
                table_team.cell(0, c_idx).text = str(col_name).replace('% Achievement from Target', 'Target %')
                table_team.cell(0, c_idx).fill.solid()
                table_team.cell(0, c_idx).fill.fore_color.rgb = RGBColor(15, 23, 42)
                if table_team.cell(0, c_idx).text_frame.paragraphs:
                    table_team.cell(0, c_idx).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    table_team.cell(0, c_idx).text_frame.paragraphs[0].font.bold = True
            
            for r_idx, row in sc_g.reset_index(drop=True).iterrows():
                for c_idx, col_name in enumerate(cols_to_show):
                    table_team.cell(r_idx + 1, c_idx).text = str(row[col_name])
                    if table_team.cell(r_idx + 1, c_idx).text_frame.paragraphs:
                        table_team.cell(r_idx + 1, c_idx).text_frame.paragraphs[0].font.size = Pt(12)

        # --- SLIDE 4, 5, 6: Dynamic Charts ---
        add_plot_to_slide(prs, fig_d, "Daily Volume & Schedule Workload Analysis")
        add_plot_to_slide(prs, fig_r, "Time Segmentation: Hourly Flow & Avg First Response (FRT)")
        add_plot_to_slide(prs, fig_st, "Time Segmentation: Hourly Flow & Avg Service Time (TAT)")

        # --- SLIDE 7: Top Request Types (Chart format) ---
        if req_counts_dict:
            types_df = pd.DataFrame(list(req_counts_dict.items()), columns=["Request Type", "Count"]).head(8)
            fig_req = px.bar(types_df, x="Request Type", y="Count", text="Count", title="Top Request Types Breakdown", color_discrete_sequence=["#3b82f6"])
            fig_req.update_traces(textposition="outside")
            fig_req.update_layout(template="plotly_white", margin=dict(l=10, r=10, t=40, b=10))
            add_plot_to_slide(prs, fig_req, "Top Request Types Breakdown")

        # --- SLIDE 8: HIC Distribution ---
        add_plot_to_slide(prs, fig_hic, "Health Insurance Companies (HIC) Distribution")

        ppt_stream = BytesIO()
        prs.save(ppt_stream)
        ppt_stream.seek(0)
        return ppt_stream, None
    except ImportError:
        return None, "المكتبة الخاصة بالباوربوينت غير متوفرة. يرجى إضافة 'python-pptx' إلى ملف requirements.txt"

# ══════════════════════════════════════════════════════════════════════════════════
#  DASHBOARD MAIN MODULE
# ══════════════════════════════════════════════════════════════════════════════════
period_ovs = overrides().get(PERIOD_KEY, {})
global_target = float(period_ovs.get("GLOBAL_TARGET", 0))

caption_text = f"Search Period: {d_from} ({DAYS_AR.get(pd.to_datetime(d_from).day_name(), '')})" if d_from == d_to else f"Search Period: {d_from} to {d_to}"
st.markdown("## 💊 In-Store Requests Matrix")
st.caption(caption_text)

# ══════════════════════════════════════════════════════════════════════════════════
#  TABS NAVIGATION ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════════
tabs_list = ["📈 Operational Insights"]
if is_admin(): tabs_list.extend(["👥 Team Performance & KPIs", "✏️ Manual Overrides", "🧑‍🤝‍🧑 Team Profiles"])
elif st.session_state.role == "expert": tabs_list.extend(["👥 Team Performance & KPIs", "🧑‍🤝‍🧑 Team Profiles"])

tabs = st.tabs(tabs_list)

# ── TAB 1 — Operational Insights (Visible to Everyone) ────────────────────────────
with tabs[0]:
    c1, c2 = st.columns(2)
    with c1: esc  = st.checkbox("🔥 Escalated Cases Only",    value=False, key="t1_esc")
    with c2: nesc = st.checkbox("🟢 Non-Escalated Cases Only", value=False, key="t1_nesc")

    dfm = df.copy(); dfm_prev = df_prev_all.copy()
    if esc and not nesc: dfm = dfm[dfm["Is Email"] == True]; dfm_prev = dfm_prev[dfm_prev["Is Email"] == True]
    elif nesc and not esc: dfm = dfm[dfm["Is Email"] == False]; dfm_prev = dfm_prev[dfm_prev["Is Email"] == False]

    global_period_adjs = overrides().get(PERIOD_KEY, {}).get("agent_adjustments", {})
    global_adj_total = sum(int(v) for v in global_period_adjs.values()) if not esc and not nesc else 0

    global_prev_adjs = overrides().get(f"{prev_d_from}_{prev_d_to}", {}).get("agent_adjustments", {})
    global_prev_adj_total = sum(int(v) for v in global_prev_adjs.values()) if not esc and not nesc else 0

    total = len(dfm) + global_adj_total
    ss    = dfm["Status"].astype(str).str.strip()
    ok    = dfm[ss.str.contains("Closed", na=False, case=False) & ~ss.str.contains("issue", na=False, case=False)].shape[0] + global_adj_total
    issue = dfm[ss.str.contains("Closed", na=False, case=False) & ss.str.contains("issue", na=False, case=False)].shape[0]
    
    curr_afr_val = dfm["Response Take (min)"].mean() if "Response Take (min)" in dfm.columns and not dfm.empty else 0
    curr_tat_val = dfm["Request Take (min)"].mean() if "Request Take (min)" in dfm.columns and not dfm.empty else 0
    
    ok_pct = (ok / total * 100) if total > 0 else 0
    issue_pct = (issue / total * 100) if total > 0 else 0
    stores_count = dfm[dfm["Store ID"] != "Unknown"]["Store ID"].nunique() if not dfm.empty else 0
    status_actions_sum = int(dfm["Status Count"].sum()) if not dfm.empty else 0
    curr_avg_per_day = (total + global_jhah + global_support) / delta_days if delta_days > 0 else 0
    
    prev_total = len(dfm_prev) + global_prev_adj_total
    ss_prev    = dfm_prev["Status"].astype(str).str.strip()
    prev_ok    = dfm_prev[ss_prev.str.contains("Closed", na=False, case=False) & ~ss_prev.str.contains("issue", na=False, case=False)].shape[0] + global_prev_adj_total
    prev_issue = dfm_prev[ss_prev.str.contains("Closed", na=False, case=False) & ss_prev.str.contains("issue", na=False, case=False)].shape[0]
    
    prev_afr_val = dfm_prev["Response Take (min)"].mean() if "Response Take (min)" in dfm_prev.columns and not dfm_prev.empty else 0
    prev_tat_val = dfm_prev["Request Take (min)"].mean() if "Request Take (min)" in dfm_prev.columns and not dfm_prev.empty else 0
    
    prev_ok_pct = (prev_ok / prev_total * 100) if prev_total > 0 else 0
    prev_issue_pct = (prev_issue / prev_total * 100) if prev_total > 0 else 0
    prev_stores_count = dfm_prev[dfm_prev["Store ID"] != "Unknown"]["Store ID"].nunique() if not dfm_prev.empty else 0
    prev_status_actions_sum = int(dfm_prev["Status Count"].sum()) if not dfm_prev.empty else 0
    prev_avg_per_day = (prev_total + prev_global_jhah + prev_global_support) / delta_days if delta_days > 0 else 0

    chg_total = calc_change(total, prev_total)
    chg_stores = calc_change(stores_count, prev_stores_count)
    chg_actions = calc_change(status_actions_sum, prev_status_actions_sum)
    chg_ok = ok_pct - prev_ok_pct  
    chg_issue = issue_pct - prev_issue_pct 
    chg_avg_per_day = calc_change(curr_avg_per_day, prev_avg_per_day)
    chg_afr = calc_change(curr_afr_val, prev_afr_val)
    chg_tat = calc_change(curr_tat_val, prev_tat_val)

    r1c1, r1c2, r1c3, r1c4, r1c5 = st.columns(5)
    r2c1, r2c2, r2c3, r2c4, r2c5 = st.columns([1.4, 0.9, 0.9, 0.9, 0.9])
    
    r1c1.markdown(kpi_colored("Total Requests",      f"{total:,}", "card-primary", chg_total, neutral=True),     unsafe_allow_html=True)
    r1c2.markdown(kpi_colored("Stores Served",      f"{stores_count:,}", "card-neutral", chg_stores, neutral=True),  unsafe_allow_html=True)
    r1c3.markdown(kpi_colored("Total Actions",      f"{status_actions_sum:,}", "card-neutral", chg_actions, neutral=True),  unsafe_allow_html=True)
    r1c4.markdown(kpi_colored("Closed Completed",   f"{ok:,} <span style='font-size:1.15rem; opacity:0.8;'>({ok_pct:.1f}%)</span>",    "card-success", chg_ok), unsafe_allow_html=True)
    r1c5.markdown(kpi_colored("Closed with Issue", f"{issue:,} <span style='font-size:1.15rem; opacity:0.8;'>({issue_pct:.1f}%)</span>", "card-danger", chg_issue, inverse=True),     unsafe_allow_html=True)
    
    r2c1.markdown(kpi_colored("Avg Requests / Day",  f"{curr_avg_per_day:.1f}", "card-neutral card-small", chg_avg_per_day, neutral=True), unsafe_allow_html=True)
    r2c2.markdown(kpi_colored("AFR (Avg Response)", fmt_m(curr_afr_val), "card-neutral card-small", chg_afr, inverse=True),       unsafe_allow_html=True)
    r2c3.markdown(kpi_colored("Avg Service (TAT)", fmt_m(curr_tat_val),        "card-neutral card-small", chg_tat, inverse=True),       unsafe_allow_html=True)
    r2c4.markdown(kpi_colored("JHAH Requests", f"{global_jhah:,}", "card-neutral card-small", neutral=True), unsafe_allow_html=True)
    r2c5.markdown(kpi_colored("Support Requests", f"{global_support:,}", "card-neutral card-small", neutral=True), unsafe_allow_html=True)
    st.write("")

    req_counts, req_pct = pd.Series(dtype=int), pd.Series(dtype=float)
    if not dfm.empty:
        req_counts = dfm['Request Type'].value_counts()
        req_pct = (req_counts / len(dfm) * 100).round(1)
        dfm_work = dfm.copy()
        dfm_work["Response Tier"] = dfm_work["Response Take (min)"].apply(assign_time_tier)
        dfm_work["Service Tier"]  = dfm_work["Request Take (min)"].apply(assign_time_tier)
        bar_palette = ["#1e40af", "#3b82f6", "#0ea5e9", "#0284c7", "#0d9488", "#14b8a6", "#475569", "#64748b", "#8b5cf6", "#a855f7"]
        bar_labels, bar_values, bar_counts = list(req_pct.index), list(req_pct.values), [int(x) for x in req_counts.values]
        bar_colors  = [bar_palette[i % len(bar_palette)] for i in range(len(req_pct))]
        SB_COLORS = {"Response Time": "#3b82f6", "Service Resolution": "#0ea5e9", "Under 15 Mins": "#22c55e", "15-30 Mins": "#3b82f6", "30-45 Mins": "#8b5cf6", "45-60 Mins": "#f59e0b", "Over 1 Hour": "#ef4444"}
        tier_order = ["Under 15 Mins", "15-30 Mins", "30-45 Mins", "45-60 Mins", "Over 1 Hour"]
        sb_payloads = {}   
        
        for rt in ["All Types"] + bar_labels:
            if rt == "All Types": sub, use_response = dfm_work.copy(), True
            else: sub, use_response = dfm_work[dfm_work["Request Type"] == rt].copy(), False
            if sub.empty: sb_payloads[rt] = {"ids": [], "labels": [], "parents": [], "values": [], "colors": []}; continue
            ids, lbl, par, val, col = [], [], [], [], []
            if use_response:
                rd, sd = sub.groupby("Response Tier").size(), sub.groupby("Service Tier").size()
                if rd.sum() > 0:
                    ids.append("Response Time"); lbl.append("Response Time"); par.append(""); val.append(int(rd.sum())); col.append(SB_COLORS["Response Time"])
                    for tier in tier_order:
                        if rd.get(tier, 0) > 0: ids.append(f"RT_{tier}"); lbl.append(tier); par.append("Response Time"); val.append(int(rd.get(tier, 0))); col.append(SB_COLORS.get(tier, "#cbd5e1"))
                if sd.sum() > 0:
                    ids.append("Service Resolution"); lbl.append("Service Resolution"); par.append(""); val.append(int(sd.sum())); col.append(SB_COLORS["Service Resolution"])
                    for tier in tier_order:
                        if sd.get(tier, 0) > 0: ids.append(f"SR_{tier}"); lbl.append(tier); par.append("Service Resolution"); val.append(int(sd.get(tier, 0))); col.append(SB_COLORS.get(tier, "#cbd5e1"))
            else:
                sd = sub.groupby("Service Tier").size()
                if sd.sum() > 0:
                    ids.append("Service Resolution"); lbl.append("Service Resolution"); par.append(""); val.append(int(sd.sum())); col.append(SB_COLORS["Service Resolution"])
                    for tier in tier_order:
                        if sd.get(tier, 0) > 0: ids.append(f"SR_{tier}"); lbl.append(tier); par.append("Service Resolution"); val.append(int(sd.get(tier, 0))); col.append(SB_COLORS.get(tier, "#cbd5e1"))
            sb_payloads[rt] = {"ids": ids, "labels": lbl, "parents": par, "values": val, "colors": col}

        import json as _json
        sb_payloads_json, bar_data_json = _json.dumps(sb_payloads), _json.dumps({"labels": bar_labels, "values": bar_values, "counts": bar_counts, "colors": bar_colors})
        component_height = max(650, len(bar_labels) * 55 + 100)

        html_component = f"""
<!DOCTYPE html><html><head><meta charset="utf-8"><script src="https://cdn.plot.ly/plotly-2.35.2.min.js"></script>
<style>* {{ box-sizing: border-box; margin: 0; padding: 0; }} body {{ background: transparent; font-family: inherit; }}
#wrapper {{ display: flex; gap: 20px; width: 100%; height: {component_height}px; align-items: stretch; }}
.chart-card {{ background: #ffffff; border-radius: 12px; padding: 20px; box-shadow: 0 1px 3px rgba(0,0,0,0.05), 0 1px 2px rgba(0,0,0,0.03); border: 1px solid #e2e8f0; display: flex; flex-direction: column; }}
#sunburst-col {{ flex: 5.5; }} #bar-col {{ flex: 4.5; }}
#sb-title {{ font-size: 1.25rem; font-weight: 800; color: #1e293b; margin-bottom: 6px; min-height: 40px; line-height: 1.3; }}
#sb-title span {{ color: #2563eb; font-size: 1.05rem; }}
#bar-title {{ font-size: 1.05rem; font-weight: 700; color: #1e293b; margin-bottom: 4px; }}
#bar-hint {{ font-size: 0.85rem; color: #64748b; margin-bottom: 6px; }}
#sunburst-div, #bar-div {{ flex: 1; min-height: 0; }}
</style></head><body>
<div id="wrapper">
  <div id="sunburst-col" class="chart-card"><div id="sb-title">&#x23F1;&#xFE0F; Service Time Breakdown (AFR &amp; TAT)</div><div id="sunburst-div"></div></div>
  <div id="bar-col" class="chart-card"><div id="bar-title">&#x1F39B;&#xFE0F; Interactive Request Types</div><div id="bar-hint">&#x1F5B1;&#xFE0F; Click any bar to filter &bull; Click empty space to reset</div><div id="bar-div"></div></div>
</div>
<script>
const SB_PAYLOADS = {sb_payloads_json}; const BAR_DATA = {bar_data_json};
let selectedRt = "All Types"; let selectedBarIdx = null;
const barTrace = {{ type: "bar", orientation: "h", x: BAR_DATA.values, y: BAR_DATA.labels, customdata: BAR_DATA.counts, text: BAR_DATA.values.map((v, i) => BAR_DATA.counts[i] + " (" + v.toFixed(1) + "%)"), textposition: "inside", insidetextanchor: "middle", textfont: {{ color: "#ffffff", size: 13, weight: "bold" }}, marker: {{ color: BAR_DATA.colors, opacity: BAR_DATA.colors.map(() => 1), line: {{ color: 'rgba(0,0,0,0.1)', width: 1 }} }}, hovertemplate: "<b>%{{y}}</b><br>Tickets: %{{customdata}}<br>Share: %{{x:.1f}}%<extra></extra>" }};
const barLayout = {{ margin: {{ l: 8, r: 8, t: 10, b: 10 }}, bargap: 0.3, xaxis: {{ visible: false }}, yaxis: {{ autorange: "reversed", tickfont: {{ size: 12, color: "#334155", weight: "bold" }}, fixedrange: true, automargin: true }}, plot_bgcolor: "rgba(0,0,0,0)", paper_bgcolor: "rgba(0,0,0,0)", font: {{ color: "#1e293b" }}, autosize: true }};
Plotly.newPlot("bar-div", [barTrace], barLayout, {{ displayModeBar: false, responsive: true }});
function buildTrace(rt) {{
  const d = SB_PAYLOADS[rt] || SB_PAYLOADS["All Types"];
  if (!d || d.labels.length === 0) return null;
  return {{ type: "sunburst", ids: d.ids, labels: d.labels, parents: d.parents, values: d.values, branchvalues: "total", sort: false, marker: {{ colors: d.colors }}, texttemplate: d.labels.map((lbl, i) => d.parents[i] === "" ? "<b>%{{label}}</b>" : "%{{label}}<br>%{{percentParent:.0%}}"), textinfo: "none", insidetextorientation: "radial", hovertemplate: "<b>%{{label}}</b><br>Tickets: %{{value:,}}<br>Share: %{{percentParent:.1%}}<extra></extra>", leaf: {{ opacity: 0.93 }} }};
}}
const sbLayout = {{ margin: {{ t: 10, b: 10, l: 10, r: 10 }}, paper_bgcolor: "rgba(0,0,0,0)", plot_bgcolor: "rgba(0,0,0,0)", font: {{ color: "#1e293b" }}, autosize: true }};
const allFrames = Object.keys(SB_PAYLOADS).map(rt => ({{ name: rt, data: [buildTrace(rt)] }})).filter(f => f.data[0] !== null);
Plotly.newPlot("sunburst-div", [buildTrace("All Types")], sbLayout, {{ displayModeBar: false, responsive: true }}).then(() => {{ Plotly.addFrames("sunburst-div", allFrames); }});
document.getElementById("bar-div").on("plotly_click", function(data) {{
  const clicked = data.points[0].y; const clickedIdx = data.points[0].pointIndex;
  if (selectedRt === clicked) {{ selectedRt = "All Types"; selectedBarIdx = null; }} else {{ selectedRt = clicked; selectedBarIdx = clickedIdx; }}
  Plotly.animate("sunburst-div", [selectedRt], {{ transition: {{ duration: 600, easing: "cubic-in-out" }}, frame: {{ duration: 600, redraw: true }} }});
  Plotly.restyle("bar-div", {{ "marker.opacity": [BAR_DATA.labels.map((_, i) => selectedRt === "All Types" ? 1 : (i === selectedBarIdx ? 1 : 0.25))] }});
  document.getElementById("sb-title").innerHTML = selectedRt === "All Types" ? "&#x23F1;&#xFE0F; Service Time Breakdown (AFR &amp; TAT)" : "&#x23F1;&#xFE0F; Service Time Breakdown (TAT)<br><span>&#x27A4; " + selectedRt + "</span>";
}});
document.getElementById("bar-div").on("plotly_deselect", function() {{
  if (selectedRt !== "All Types") {{
    selectedRt = "All Types"; selectedBarIdx = null;
    Plotly.animate("sunburst-div", [selectedRt], {{ transition: {{ duration: 600, easing: "cubic-in-out" }}, frame: {{ duration: 600, redraw: true }} }});
    Plotly.restyle("bar-div", {{ "marker.opacity": [BAR_DATA.labels.map(() => 1)] }});
    document.getElementById("sb-title").innerHTML = "&#x23F1;&#xFE0F; Service Time Breakdown (AFR &amp; TAT)";
  }}
}});
</script></body></html>
"""
        import streamlit.components.v1 as components
        components.html(html_component, height=component_height + 40, scrolling=False)

    st.divider()

    fig_r, fig_st, fig_d, fig_hic = None, None, None, None

    if not df_raw.empty:
        st.markdown("### ⏳ Ticket flow rate over daily hours")
        df_flow_strict = df_raw[(df_raw["Date Only"] >= d_from) & (df_raw["Date Only"] <= d_to)].copy()
        if esc and not nesc: df_flow_strict = df_flow_strict[df_flow_strict["Is Email"] == True]
        elif nesc and not esc: df_flow_strict = df_flow_strict[df_flow_strict["Is Email"] == False]
            
        hrs = df_flow_strict.groupby("Hour").agg(
            Volume=("Request ID", "count"), 
            AR=("Response Take (min)" , "mean"),
            ST=("Request Take (min)", "mean")
        ).reset_index().set_index("Hour").reindex(range(24)).fillna(0).reset_index()
        
        hrs["Hour Label"] = ["12 AM" if h == 0 else ("12 PM" if h == 12 else (f"{h} AM" if h < 12 else f"{h - 12} PM")) for h in hrs["Hour"]]
        hrs["Avg_Vol"] = hrs["Volume"] / max(1, delta_days)
        
        fig_r = make_subplots(specs=[[{"secondary_y": True}]])
        fig_r.add_trace(go.Scatter(x=hrs["Hour Label"], y=hrs["Volume"], name="Volume", fill="tozeroy", line=dict(color="#3b82f6", width=2), customdata=hrs["Avg_Vol"], hovertemplate="%{y} (Avg: %{customdata:.1f}/day)<extra></extra>"), secondary_y=False)
        fig_r.add_trace(go.Scatter(x=hrs["Hour Label"], y=hrs["AR"], name="FRT (Avg Response)", mode="lines+markers", line=dict(color="#10b981", width=3, shape="spline"), hovertemplate="%{y:.1f} min<extra></extra>"), secondary_y=True)
        fig_r.update_layout(template="plotly_white", paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", font_color="#1e293b", margin=dict(l=10, r=10, t=55, b=10), height=550, hovermode="x unified", legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1))
        fig_r.update_yaxes(title_text="Volume (Tickets)", secondary_y=False); fig_r.update_yaxes(title_text="Avg Response Time (min)", secondary_y=True)
        st.plotly_chart(fig_r, use_container_width=True)

        st.divider()

        st.markdown("### ⏱️ Service Time (TAT) over daily hours")
        fig_st = make_subplots(specs=[[{"secondary_y": True}]])
        fig_st.add_trace(go.Scatter(x=hrs["Hour Label"], y=hrs["Volume"], name="Volume", fill="tozeroy", line=dict(color="#3b82f6", width=2), customdata=hrs["Avg_Vol"], hovertemplate="%{y} (Avg: %{customdata:.1f}/day)<extra></extra>"), secondary_y=False)
        fig_st.add_trace(go.Scatter(x=hrs["Hour Label"], y=hrs["ST"], name="TAT (Avg Service)", mode="lines+markers", line=dict(color="#f59e0b", width=3, shape="spline"), hovertemplate="%{y:.1f} min<extra></extra>"), secondary_y=True)
        fig_st.update_layout(template="plotly_white", paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", font_color="#1e293b", margin=dict(l=10, r=10, t=55, b=10), height=550, hovermode="x unified", legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1))
        fig_st.update_yaxes(title_text="Volume (Tickets)", secondary_y=False); fig_st.update_yaxes(title_text="Avg Service Time (min)", secondary_y=True)
        st.plotly_chart(fig_st, use_container_width=True)

        st.divider()
        st.markdown("### 📅 Daily Volume & Schedule Workload Analysis")
        st.markdown("""<div style="text-align: left; font-size: 1.1rem; margin-bottom: 1rem; color: #475569;"><strong>Agent Workload Indicator (Tickets per Agent):</strong><br><span style="display: inline-block; margin-right: 15px;"><span style="display:inline-block; width:14px; height:14px; background-color:#3b82f6; border-radius:3px; vertical-align:middle; margin-right:6px; margin-bottom:2px;"></span>Optimal (≤55)</span><span style="display: inline-block; margin-right: 15px;"><span style="display:inline-block; width:14px; height:14px; background-color:#eab308; border-radius:3px; vertical-align:middle; margin-right:6px; margin-bottom:2px;"></span>Moderate (56-60)</span><span style="display: inline-block; margin-right: 15px;"><span style="display:inline-block; width:14px; height:14px; background-color:#f97316; border-radius:3px; vertical-align:middle; margin-right:6px; margin-bottom:2px;"></span>High (61-63)</span><span style="display: inline-block; margin-right: 15px;"><span style="display:inline-block; width:14px; height:14px; background-color:#ef4444; border-radius:3px; vertical-align:middle; margin-right:6px; margin-bottom:2px;"></span>Severe (64-70)</span><span style="display: inline-block;"><span style="display:inline-block; width:14px; height:14px; background-color:#991b1b; border-radius:3px; vertical-align:middle; margin-right:6px; margin-bottom:2px;"></span>Excessive (>70)</span></div>""", unsafe_allow_html=True)
        
        df_workload = df_raw[(df_raw["Date Only"] >= d_from) & (df_raw["Date Only"] <= d_to)].copy()
        if esc and not nesc: df_workload = df_workload[df_workload["Is Email"] == True]
        elif nesc and not esc: df_workload = df_workload[df_workload["Is Email"] == False]
        dfm_shift = df_workload.copy(); dfm_shift["Shift Date"] = dfm_shift["Date Only"]
        daily_vol = dfm_shift.groupby("Shift Date").agg(Total_Tickets=("Request ID", "count")).reset_index()
        
        roster_date_map = {}
        if not df_roster.empty:
            for col in df_roster.columns:
                match = re.search(r'(\d{1,2}[-/\s]+(?:[A-Za-z]+|\d{1,2})[-/\s]+\d{4}|\d{4}[-/\s]+\d{1,2}[-/\s]+\d{1,2})', re.sub(r'(?i)\b(monday|tuesday|wednesday|thursday|friday|saturday|sunday)\b', '', str(col).strip()).strip())
                if match:
                    try: roster_date_map[pd.to_datetime(match.group(1), dayfirst=True).date()] = col
                    except: pass
        
        tracked_ids = [str(v).strip().lower() for v in EXPERT_ID_MAP.values()]
        def get_scheduled_agents(target_date):
            if target_date not in roster_date_map or df_roster.empty: return -1 
            col_name = roster_date_map[target_date]
            working_count = 0
            for _, row in df_roster.iterrows():
                if any(tid in " ".join([str(x).strip().lower() for x in row.values]) for tid in tracked_ids):
                    cell_val = str(row.get(col_name, "")).strip().lower()
                    if cell_val and cell_val not in EXCLUSION_LIST: working_count += 1
            return working_count
        
        daily_vol["Scheduled_Agents"] = daily_vol["Shift Date"].apply(get_scheduled_agents)
        active_df = dfm_shift[dfm_shift["Assigned By"].isin(OFFICIAL_EXPERTS)].groupby("Shift Date").agg(Actual_Agents=("Assigned By", "nunique")).reset_index()
        daily_vol = pd.merge(daily_vol, active_df, on="Shift Date", how="left")
        daily_vol["Actual_Agents"] = daily_vol["Actual_Agents"].fillna(0)
        daily_vol["Active_Agents"] = np.where(daily_vol["Scheduled_Agents"] != -1, daily_vol["Scheduled_Agents"], daily_vol["Actual_Agents"]).clip(1)
        daily_vol["Tickets per Agent"] = (daily_vol["Total_Tickets"] / daily_vol["Active_Agents"]).round(1)
        daily_vol["Date DT"] = pd.to_datetime(daily_vol["Shift Date"])
        daily_vol["Day Name"] = daily_vol["Date DT"].dt.day_name()
        DAY_COLORS = {"Saturday": "#64748b", "Sunday": "#1e40af", "Monday": "#2563eb", "Tuesday": "#3b82f6", "Wednesday": "#0ea5e9", "Thursday": "#0284c7", "Friday": "#475569"}
        daily_vol["Date Label"] = daily_vol.apply(lambda r: f"{r['Date DT'].strftime('%b %d')}<br><span style='color:{DAY_COLORS.get(r['Day Name'], '#1e293b')}'><b>({r['Day Name']})</b></span>", axis=1)
        daily_vol["Color"] = np.select([daily_vol["Tickets per Agent"] > 70, daily_vol["Tickets per Agent"] > 63, daily_vol["Tickets per Agent"] > 60, daily_vol["Tickets per Agent"] > 55], ["#991b1b", "#ef4444", "#f97316", "#eab308"], default="#3b82f6") 
        
        fig_d = make_subplots(specs=[[{"secondary_y": True}]])
        fig_d.add_trace(go.Bar(x=daily_vol["Date Label"], y=daily_vol["Total_Tickets"], text=daily_vol["Total_Tickets"], textposition='auto', marker_color=daily_vol["Color"], name="Total Tickets", showlegend=False, hovertemplate="<b>%{x}</b><br>Tickets: %{y}<extra></extra>"), secondary_y=False)
        fig_d.add_trace(go.Scatter(x=daily_vol["Date Label"], y=daily_vol["Tickets per Agent"], name="Tickets per Agent (Workload)", mode="lines+markers+text", text=daily_vol["Tickets per Agent"], textposition="top center", line=dict(color="#475569", width=3, shape="spline"), marker=dict(size=8, color="#1e293b"), hovertemplate="<b>%{x}</b><br>Tickets/Agent: %{y}<br>Scheduled Agents: %{customdata}<extra></extra>", customdata=daily_vol["Active_Agents"]), secondary_y=True)
        fig_d.update_layout(template="plotly_white", paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", font_color="#1e293b", margin=dict(l=10, r=10, t=55, b=10), height=480, xaxis_title="", hovermode="x unified", legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1))
        fig_d.update_yaxes(title_text="Total Tickets Count", secondary_y=False); fig_d.update_yaxes(title_text="Workload Ratio (Tickets/Agent)", secondary_y=True)
        st.plotly_chart(fig_d, use_container_width=True)

        st.divider()
        st.markdown("### 🏥 Health Insurance Companies (HIC) Distribution Analysis")
        df_hic_strict = df_raw[(df_raw["Date Only"] >= d_from) & (df_raw["Date Only"] <= d_to)].copy()
        if esc and not nesc: df_hic_strict = df_hic_strict[df_hic_strict["Is Email"] == True]
        elif nesc and not esc: df_hic_strict = df_hic_strict[df_hic_strict["Is Email"] == False]
        if not df_hic_strict.empty:
            hic_counts = df_hic_strict.groupby("HIC").agg(Volume=("Request ID", "count")).reset_index().sort_values(by="Volume", ascending=False) 
            fig_hic = px.bar(hic_counts, x="HIC", y="Volume", text="Volume", color_discrete_sequence=["#2563eb"], labels={"Volume": "Tickets Count", "HIC": "Insurance Provider"})
            fig_hic.update_traces(textposition="outside", hovertemplate="<b>%{x}</b><br>Tickets Resolved: %{y:,}<extra></extra>")
            fig_hic.update_layout(template="plotly_white", paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", font_color="#1e293b", margin=dict(l=10, r=10, t=55, b=10), height=500, xaxis_title="", yaxis_title="Total Handled Volume (Tickets)", xaxis_tickangle=-45, legend_title_text="Insurance Provider")
            st.plotly_chart(fig_hic, use_container_width=True)
        else: st.info("No insurance (HIC) records available for this period.")

    if is_admin():
        st.divider(); st.markdown("#### 📥 Export Operational Report")
        c_exp1, c_exp2 = st.columns(2)
        with c_exp1:
            base_metrics = [{"Metric": "Total Tickets", "Value": total}, {"Metric": "Stores Served", "Value": stores_count}, {"Metric": "Total Actions", "Value": status_actions_sum}, {"Metric": "Closed Completed", "Value": ok}, {"Metric": "Closed with Issue", "Value": issue}, {"Metric": "Avg Requests / Day", "Value": round(curr_avg_per_day, 1)}, {"Metric": "AFR", "Value": fmt_m(curr_afr_val)}, {"Metric": "TAT", "Value": fmt_m(curr_tat_val)}, {"Metric": "JHAH Requests", "Value": global_jhah}, {"Metric": "Support Requests", "Value": global_support}]
            if not dfm.empty:
                base_metrics.append({"Metric": "--- AVG TICKETS PER WEEKDAY ---", "Value": ""})
                avg_per_weekday = dfm.groupby(['Date Only', 'Day Name']).size().reset_index(name='Tickets').groupby('Day Name')['Tickets'].mean().round(1)
                for d in ['Saturday', 'Sunday', 'Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday']: base_metrics.append({"Metric": f"Avg Tickets ({d})", "Value": avg_per_weekday.get(d, 0)})
                base_metrics.append({"Metric": "--- REQUEST TYPES BREAKDOWN ---", "Value": ""})
                has_tat = "Request Take (min)" in dfm.columns
                rt_tat = dfm.groupby("Request Type")["Request Take (min)"].mean() if has_tat else pd.Series(dtype=float)
                for rt_name, rt_count in req_counts.items(): base_metrics.append({"Metric": f"Type: {rt_name}", "Value": f"{rt_count} Tickets ({req_pct[rt_name]}%) | Avg TAT: {fmt_m(rt_tat.get(rt_name, 0)) if has_tat else '00:00:00'}"})
            st.download_button("📥 Download Operational Summary (CSV)", pd.DataFrame(base_metrics).to_csv(index=False).encode('utf-8-sig'), f"Operational_Summary_{d_from}_to_{d_to}.csv", "text/csv", use_container_width=True)
        with c_exp2:
            # === تحضير بيانات الفريق أوتوماتيك لبريزنتيشن الباوربوينت ===
            df_sc_g = df[df["Assigned By"].astype(str).str.strip().str.lower().isin([x.lower() for x in OFFICIAL_EXPERTS])].copy()
            sc_g = pd.DataFrame({"Expert": OFFICIAL_EXPERTS})
            if not df_sc_g.empty:
                df_sc_g["Assigned By"] = df_sc_g["Assigned By"].astype(str).str.strip().str.lower().map({x.lower(): x for x in OFFICIAL_EXPERTS})
                grp_g = df_sc_g.groupby("Assigned By")
                stats_g = pd.DataFrame(index=grp_g.groups.keys())
                stats_g["Tickets Count"] = grp_g["Request ID"].count()
                sc_g = sc_g.merge(stats_g, left_on="Expert", right_index=True, how="left")
            else:
                sc_g["Tickets Count"] = 0
                
            sc_g["Tickets Count"] = sc_g["Tickets Count"].fillna(0).astype(int)
            period_adjs_g = overrides().get(PERIOD_KEY, {}).get("agent_adjustments", {})
            sc_g["Tickets Count"] = sc_g.apply(lambda row: row["Tickets Count"] + int(period_adjs_g.get(str(row["Expert"]).strip(), 0)), axis=1)
            sc_g["Working Days"] = sc_g["Expert"].apply(lambda x: curr_roster_counts.get(x, {}).get("wd", 0))
            sc_g["JHAH Requests"] = sc_g.apply(lambda row: int(float(out_req_dict.get(str(row["Expert"]).strip().lower(), {}).get("JHAH", 0) or 0)), axis=1)
            sc_g["Support Requests"] = sc_g.apply(lambda row: int(float(out_req_dict.get(str(row["Expert"]).strip().lower(), {}).get("Support Req", 0) or 0)), axis=1)
            
            total_cases_g = pd.to_numeric(sc_g["Tickets Count"], errors='coerce').fillna(0) + pd.to_numeric(sc_g["JHAH Requests"], errors='coerce').fillna(0) + pd.to_numeric(sc_g["Support Requests"], errors='coerce').fillna(0)
            sc_g["Cases/Day"] = (total_cases_g / pd.to_numeric(sc_g["Working Days"], errors='coerce').fillna(0).replace(0, 1)).round(1)
            
            if global_target > 0: sc_g["% Achievement from Target"] = ((sc_g["Cases/Day"] / global_target) * 100).round(1).astype(str) + "%"
            else: sc_g["% Achievement from Target"] = ((sc_g["Cases/Day"] / sc_g["Cases/Day"].mean() * 100).round(1).astype(str) + "%" if sc_g["Cases/Day"].mean() > 0 else "0.0%")
            
            sc_g["Service Quality"] = sc_g["Expert"].apply(lambda x: f"{100.0 - float(expert_quality_deductions.get(str(x).strip().lower(), 0.0)):.1f}%")
            
            if not sc_g.empty:
                sc_g["_sort_qual"] = sc_g["Service Quality"].astype(str).str.replace('%', '', regex=False).astype(float)
                sc_g["_sort_cases"] = sc_g["Cases/Day"].astype(float)
                sc_g["_rank_score"] = (sc_g["_sort_qual"] * 1000) + sc_g["_sort_cases"]
                sc_g.sort_values(by="_rank_score", ascending=False, inplace=True)
                sc_g.insert(1, "Rank", sc_g["_rank_score"].rank(method="min", ascending=False).astype(int).astype(str))
            else:
                sc_g.insert(1, "Rank", [])
                
            ppt_file_data, err_msg = generate_pptx_summary(d_from.strftime('%Y-%m-%d'), d_to.strftime('%Y-%m-%d'), total, curr_avg_per_day, ok, ok_pct, issue, issue_pct, fmt_m(curr_afr_val), fmt_m(curr_tat_val), stores_count, status_actions_sum, global_jhah, global_support, req_counts.to_dict(), req_pct.to_dict(), sc_g, fig_r, fig_st, fig_d, fig_hic)
            
            if ppt_file_data:
                st.download_button("📊 Export to PowerPoint (PPTX)", data=ppt_file_data.getvalue(), file_name=f"Operational_Summary_{d_from}_to_{d_to}.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)
            else:
                st.warning(f"⚠️ {err_msg}")

# ── TAB 2 — Team Performance and KPIs (Admin & Expert) ─────────────────────────────
if len(tabs) > 1 and (is_admin() or st.session_state.role == "expert"):
    with tabs[1]:
        st.markdown("### 👥 Team Performance and KPIs")
        sel_agents_t2 = st.multiselect("Filter by Expert Name", sorted(df_raw["Assigned By"].dropna().unique()), key="t2_agents")
        aname = my_agent_name()
        is_exp = (st.session_state.role == "expert" and aname)
        
        target_agent_photo = aname if is_exp else (sel_agents_t2[0] if sel_agents_t2 and len(sel_agents_t2) == 1 else None)
        if target_agent_photo:
            p_photo = "https://cdn-icons-png.flaticon.com/512/3135/3135715.png"
            p_bio, p_role = "", "⭐ EXPERT"
            for u in users().values():
                if u.get("agent_name") == target_agent_photo or u.get("display_name") == target_agent_photo:
                    p_photo = u.get("photo", "") or p_photo
                    p_bio = u.get("bio", "")
                    p_role = "👑 ADMIN" if u.get('role') == 'admin' else ("👁️ SUPERVISOR" if u.get('role') == 'supervisor' else "⭐ EXPERT")
                    break
            
            st.markdown(f'''
            <div style="display: flex; align-items: center; background: #ffffff; padding: 1.2rem; border-radius: 12px; border: 1px solid #cbd5e1; margin-bottom: 1.5rem; box-shadow: 0 4px 6px -1px rgba(0,0,0,0.05);">
                <img src="{p_photo}" onerror="this.src='https://cdn-icons-png.flaticon.com/512/3135/3135715.png'" style="width: 85px; height: 85px; border-radius: 50%; object-fit: cover; object-position: top; border: 3px solid #2563eb; padding: 2px; margin-right: 1.5rem;">
                <div>
                    <h2 style="margin: 0 !important; padding: 0 !important; font-size: 1.6rem !important; color: #0f172a;">{target_agent_photo}</h2>
                    <span style="color: #64748b; font-weight: 700; font-size: 0.95rem;">{p_role}</span>
                    <p style="margin: 0.4rem 0 0 0 !important; color: #475569; font-style: italic; font-size: 0.95rem;">"{p_bio}"</p>
                </div>
            </div>
            ''', unsafe_allow_html=True)
        
        scope_agents = [aname] if is_exp else (sel_agents_t2 if sel_agents_t2 else OFFICIAL_EXPERTS)
        period_adjs = overrides().get(PERIOD_KEY, {}).get("agent_adjustments", {})
        total_kpi_adj = sum(int(period_adjs.get(agent, 0)) for agent in scope_agents)

        prev_period_adjs = overrides().get(f"{prev_d_from}_{prev_d_to}", {}).get("agent_adjustments", {})
        prev_total_kpi_adj = sum(int(prev_period_adjs.get(agent, 0)) for agent in scope_agents)
        
        df_kpi, df_kpi_prev = df.copy(), df_prev_all.copy()
        if is_exp: df_kpi = df_kpi[df_kpi["Assigned By"] == aname]; df_kpi_prev = df_kpi_prev[df_kpi_prev["Assigned By"] == aname]
        elif sel_agents_t2: df_kpi = df_kpi[df_kpi["Assigned By"].isin(sel_agents_t2)]; df_kpi_prev = df_kpi_prev[df_kpi_prev["Assigned By"].isin(sel_agents_t2)]
            
        total_kpi, prev_kpi_total = len(df_kpi) + total_kpi_adj, len(df_kpi_prev) + prev_total_kpi_adj
        
        kpi_ss, kpi_ss_prev = df_kpi["Status"].astype(str).str.strip(), df_kpi_prev["Status"].astype(str).str.strip()
        kpi_ok = df_kpi[kpi_ss.str.contains("Closed", na=False, case=False) & ~kpi_ss.str.contains("issue", na=False, case=False)].shape[0] + total_kpi_adj
        prev_kpi_ok = df_kpi_prev[kpi_ss_prev.str.contains("Closed", na=False, case=False) & ~kpi_ss_prev.str.contains("issue", na=False, case=False)].shape[0] + prev_total_kpi_adj
        kpi_iss = df_kpi[kpi_ss.str.contains("Closed", na=False, case=False) & kpi_ss.str.contains("issue", na=False, case=False)].shape[0]
        prev_kpi_iss = df_kpi_prev[kpi_ss_prev.str.contains("Closed", na=False, case=False) & kpi_ss_prev.str.contains("issue", na=False, case=False)].shape[0]
        
        kpi_curr_afr_val = df_kpi["Response Take (min)"].mean() if "Response Take (min)" in df_kpi.columns and not df_kpi.empty else 0
        kpi_curr_tat_val = df_kpi["Request Take (min)"].mean() if "Request Take (min)" in df_kpi.columns and not df_kpi.empty else 0
        prev_kpi_afr_val = df_kpi_prev["Response Take (min)"].mean() if "Response Take (min)" in df_kpi_prev.columns and not df_kpi_prev.empty else 0
        prev_kpi_tat_val = df_kpi_prev["Request Take (min)"].mean() if "Request Take (min)" in df_kpi_prev.columns and not df_kpi_prev.empty else 0
        
        kpi_ok_pct = (kpi_ok / total_kpi * 100) if total_kpi > 0 else 0
        kpi_iss_pct = (kpi_iss / total_kpi * 100) if total_kpi > 0 else 0
        prev_kpi_ok_pct = (prev_kpi_ok / prev_kpi_total * 100) if prev_kpi_total > 0 else 0
        prev_kpi_iss_pct = (prev_kpi_iss / prev_kpi_total * 100) if prev_kpi_total > 0 else 0

        df_sc = df[df["Assigned By"].astype(str).str.strip().str.lower().isin([x.lower() for x in OFFICIAL_EXPERTS])].copy()
        sc = pd.DataFrame({"Expert": OFFICIAL_EXPERTS})
        if not df_sc.empty:
            df_sc["Assigned By"] = df_sc["Assigned By"].astype(str).str.strip().str.lower().map({x.lower(): x for x in OFFICIAL_EXPERTS})
            df_sc["_jhah"] = df_sc["Request Type"].astype(str).str.lower().str.contains("jhah", na=False)
            df_sc["_c_ok"] = (df_sc["Status"].astype(str).str.contains("Closed", case=False, na=False) & ~df_sc["Status"].astype(str).str.contains("issue", case=False, na=False))
            df_sc["_c_all"] = df_sc["Status"].astype(str).str.contains("Closed", case=False, na=False)
            grp = df_sc.groupby("Assigned By")
            stats = pd.DataFrame(index=grp.groups.keys())
            stats["Tickets Count"] = grp["Request ID"].count()
            stats["_Service_Time_val"] = grp["Request Take (min)"].mean() if "Request Take (min)" in df_sc.columns else 0
            stats["_AFR_val"] = grp["Response Take (min)"].mean() if "Response Take (min)" in df_sc.columns else 0
            stats["_c_ok_sum"], stats["_c_all_sum"] = grp["_c_ok"].sum(), grp["_c_all"].sum()
            sc = sc.merge(stats, left_on="Expert", right_index=True, how="left")
        else: sc["Tickets Count"], sc["_Service_Time_val"], sc["_AFR_val"], sc["_c_ok_sum"], sc["_c_all_sum"] = 0, 0, 0, 0, 0
            
        sc["Tickets Count"] = sc["Tickets Count"].fillna(0).astype(int)
        
        # ── SECRET MANUAL ADJUSTMENTS ──
        sc["Tickets Count"] = sc.apply(lambda row: row["Tickets Count"] + int(period_adjs.get(str(row["Expert"]).strip(), 0)), axis=1)

        sc["Working Days"] = sc["Expert"].apply(lambda x: curr_roster_counts.get(x, {}).get("wd", 0))
        sc["Off Days"] = sc["Expert"].apply(lambda x: curr_roster_counts.get(x, {}).get("off", 0))
        sc["Annual Leaves"] = sc["Expert"].apply(lambda x: curr_roster_counts.get(x, {}).get("ann", 0))
        sc["Casual Leaves"] = sc["Expert"].apply(lambda x: curr_roster_counts.get(x, {}).get("cas", 0))
        sc["Sick Leaves"] = sc["Expert"].apply(lambda x: curr_roster_counts.get(x, {}).get("sick", 0))
        sc["AFR"], sc["Service Time"] = sc["_AFR_val"].fillna(0).apply(fmt_m), sc["_Service_Time_val"].fillna(0).apply(fmt_m)
        sc["JHAH Requests"] = sc.apply(lambda row: int(float(out_req_dict.get(str(row["Expert"]).strip().lower(), {}).get("JHAH", 0) or 0)), axis=1)
        sc["Support Requests"] = sc.apply(lambda row: int(float(out_req_dict.get(str(row["Expert"]).strip().lower(), {}).get("Support Req", 0) or 0)), axis=1)
        
        sc = sc[(pd.to_numeric(sc["Working Days"], errors='coerce').fillna(0) > 0) | ((pd.to_numeric(sc["Tickets Count"], errors='coerce').fillna(0) + pd.to_numeric(sc["JHAH Requests"], errors='coerce').fillna(0) + pd.to_numeric(sc["Support Requests"], errors='coerce').fillna(0)) > 0)].copy()
        
        total_cases = pd.to_numeric(sc["Tickets Count"], errors='coerce').fillna(0) + pd.to_numeric(sc["JHAH Requests"], errors='coerce').fillna(0) + pd.to_numeric(sc["Support Requests"], errors='coerce').fillna(0)
        sc["Cases/Day"] = (total_cases / pd.to_numeric(sc["Working Days"], errors='coerce').fillna(0).replace(0, 1)).round(1)
            
        if global_target > 0: sc["% Achievement from Target"] = ((sc["Cases/Day"] / global_target) * 100).round(1).astype(str) + "%"
        else: sc["% Achievement from Target"] = ((sc["Cases/Day"] / sc["Cases/Day"].mean() * 100).round(1).astype(str) + "%" if sc["Cases/Day"].mean() > 0 else "0.0%")
            
        sc["Service Quality"] = sc["Expert"].apply(lambda x: f"{100.0 - float(expert_quality_deductions.get(str(x).strip().lower(), 0.0)):.1f}%")
        
        def calc_incentive(row):
            try: achiev = float(str(row["% Achievement from Target"]).replace("%", "")) / 100.0
            except: achiev = 0.0
            count_inc = 1500 if achiev >= 0.97 else (1350 if achiev >= 0.95 else (1200 if achiev >= 0.90 else (1050 if achiev >= 0.85 else 0)))
            try: qual = float(str(row["Service Quality"]).replace("%", "")) / 100.0
            except: qual = 0.0
            return f"{count_inc + (600 * qual) + 1000:,.0f} EGP"
            
        sc["Prospected Incentive"] = sc.apply(calc_incentive, axis=1)

        # ── SMART RANKING SYSTEM (Incentive -> Quality -> Cases/Day) ──
        if not sc.empty:
            sc["_sort_inc"] = sc["Prospected Incentive"].astype(str).str.replace(',', '', regex=False).str.replace(' EGP', '', regex=False).astype(float)
            sc["_sort_qual"] = sc["Service Quality"].astype(str).str.replace('%', '', regex=False).astype(float)
            sc["_sort_cases"] = sc["Cases/Day"].astype(float)
            
            sc["_rank_score"] = (sc["_sort_inc"] * 100000) + (sc["_sort_qual"] * 1000) + sc["_sort_cases"]
            sc.sort_values(by="_rank_score", ascending=False, inplace=True)
            sc.insert(1, "Rank", sc["_rank_score"].rank(method="min", ascending=False).astype(int).astype(str))
            sc.drop(columns=["_sort_inc", "_sort_qual", "_sort_cases", "_rank_score"], inplace=True)
        else:
            sc.insert(1, "Rank", [])

        kpi_scope_df = sc[sc["Expert"] == aname] if is_exp else (sc[sc["Expert"].isin(sel_agents_t2)] if sel_agents_t2 else sc.copy())
        
        kpi_curr_avg_per_day = (pd.to_numeric(kpi_scope_df["Tickets Count"], errors='coerce').fillna(0).sum() + pd.to_numeric(kpi_scope_df["JHAH Requests"], errors='coerce').fillna(0).sum() + pd.to_numeric(kpi_scope_df["Support Requests"], errors='coerce').fillna(0).sum()) / pd.to_numeric(kpi_scope_df["Working Days"], errors='coerce').fillna(0).sum() if pd.to_numeric(kpi_scope_df["Working Days"], errors='coerce').fillna(0).sum() > 0 else 0

        prev_sum_wd = sum(prev_roster_counts[exp]["wd"] for exp in scope_agents)
        prev_kpi_scope_df = df_kpi_prev[df_kpi_prev["Assigned By"].str.lower().isin([a.lower() for a in scope_agents])] if scope_agents else df_kpi_prev
        prev_kpi_avg_per_day = (len(prev_kpi_scope_df) + prev_total_kpi_adj + sum(prev_out_req_dict.get(e.lower(), {}).get("JHAH", 0) for e in scope_agents) + sum(prev_out_req_dict.get(e.lower(), {}).get("Support Req", 0) for e in scope_agents)) / prev_sum_wd if prev_sum_wd > 0 else 0

        chg_kpi_total = calc_change(total_kpi, prev_kpi_total); chg_kpi_ok = kpi_ok_pct - prev_kpi_ok_pct; chg_kpi_iss = kpi_iss_pct - prev_kpi_iss_pct; chg_kpi_avg_per_day = calc_change(kpi_curr_avg_per_day, prev_kpi_avg_per_day); chg_kpi_afr = calc_change(kpi_curr_afr_val, prev_kpi_afr_val); chg_kpi_tat = calc_change(kpi_curr_tat_val, prev_kpi_tat_val)

        k1, k2, k3, k4, k5, k6 = st.columns(6)
        k1.markdown(kpi_colored("Total Requests",      f"{total_kpi:,}", "card-primary", chg_kpi_total, neutral=True),     unsafe_allow_html=True)
        k2.markdown(kpi_colored("Avg Requests / Day", f"{kpi_curr_avg_per_day:.1f}", "card-neutral", chg_kpi_avg_per_day, neutral=True),     unsafe_allow_html=True)
        k3.markdown(kpi_colored("Closed Completed",   f"{kpi_ok:,} <span style='font-size:1.15rem; opacity:0.8;'>({kpi_ok_pct:.1f}%)</span>",      "card-success", chg_kpi_ok), unsafe_allow_html=True)
        k4.markdown(kpi_colored("Closed with Issue",  f"{kpi_iss:,} <span style='font-size:1.15rem; opacity:0.8;'>({kpi_iss_pct:.1f}%)</span>",     "card-danger", chg_kpi_iss, inverse=True),     unsafe_allow_html=True)
        k5.markdown(kpi_colored("AFR (Avg First Response)", fmt_m(kpi_curr_afr_val), "card-neutral", chg_kpi_afr, inverse=True), unsafe_allow_html=True)
        k6.markdown(kpi_colored("Avg Service (TAT)",  fmt_m(kpi_curr_tat_val), "card-neutral", chg_kpi_tat, inverse=True), unsafe_allow_html=True)
        st.write(""); st.divider()

        team_wd = round(pd.to_numeric(sc["Working Days"], errors='coerce').mean(), 1) if not sc.empty else 0
        team_row = {
            "Expert": "🏆 Team AVG", "Rank": "-",
            "Working Days": team_wd, "Tickets Count": round(pd.to_numeric(sc["Tickets Count"], errors='coerce').mean(), 1) if not sc.empty else 0, "JHAH Requests": round(pd.to_numeric(sc["JHAH Requests"], errors='coerce').mean(), 1) if not sc.empty else 0, "Support Requests": round(pd.to_numeric(sc["Support Requests"], errors='coerce').mean(), 1) if not sc.empty else 0,
            "Cases/Day": round((round(pd.to_numeric(sc["Tickets Count"], errors='coerce').mean(), 1) if not sc.empty else 0 + round(pd.to_numeric(sc["JHAH Requests"], errors='coerce').mean(), 1) if not sc.empty else 0 + round(pd.to_numeric(sc["Support Requests"], errors='coerce').mean(), 1) if not sc.empty else 0) / (team_wd if team_wd > 0 else 1), 1),
            "Off Days": round(pd.to_numeric(sc["Off Days"], errors='coerce').mean(), 1) if not sc.empty else 0, "Annual Leaves": round(pd.to_numeric(sc["Annual Leaves"], errors='coerce').mean(), 1) if not sc.empty else 0, "Casual Leaves": round(pd.to_numeric(sc["Casual Leaves"], errors='coerce').mean(), 1) if not sc.empty else 0, "Sick Leaves": round(pd.to_numeric(sc["Sick Leaves"], errors='coerce').mean(), 1) if not sc.empty else 0,
            "% Achievement from Target": "100.0%", "AFR": fmt_m(df_sc["Response Take (min)"].mean() if "Response Take (min)" in df_sc.columns and not df_sc.empty else 0), "Service Time": fmt_m(df_sc["Request Take (min)"].mean() if "Request Take (min)" in df_sc.columns and not df_sc.empty else 0),
            "Service Quality": f"{sc['Service Quality'].apply(lambda p: float(str(p).replace('%', '')) if pd.notnull(p) else 0.0).mean() if not sc.empty else 0:.1f}%", "Prospected Incentive": "3,100 EGP"
        }
        sc.drop(columns=["_Service_Time_val", "_AFR_val"], inplace=True, errors='ignore')
        sc_final = pd.concat([pd.DataFrame([team_row]), sc], ignore_index=True)
        
        rank_df = sc.copy(); rank_df["_sort_val"] = pd.to_numeric(rank_df["Rank"], errors="coerce").fillna(999)
        top_exps = rank_df.nsmallest(3, "_sort_val")["Expert"].tolist()
        gold_exp, silver_exp, bronze_exp = top_exps[0] if len(top_exps) > 0 else None, top_exps[1] if len(top_exps) > 1 else None, top_exps[2] if len(top_exps) > 2 else None

        display_df = sc_final[sc_final["Expert"].isin(["🏆 Team AVG", aname])] if is_exp else (sc_final[sc_final["Expert"].isin(["🏆 Team AVG"] + sel_agents_t2)] if sel_agents_t2 else sc_final.copy())
        
        def format_clean_num(x):
            if x == "-": return "-"
            try: return str(int(float(x))) if float(x).is_integer() else str(round(float(x), 1))
            except: return str(x)
        for c in ["Working Days", "Tickets Count", "JHAH Requests", "Support Requests", "Cases/Day", "Off Days", "Annual Leaves", "Casual Leaves", "Sick Leaves"]:
            if c in display_df.columns: display_df[c] = display_df[c].apply(format_clean_num)

        display_df["Expert"] = display_df["Expert"].apply(lambda v: f"🥇 {v}" if v == gold_exp else (f"🥈 {v}" if v == silver_exp else (f"🥉 {v}" if v == bronze_exp else v)))

        def style_performers(row):
            exp, styles = row["Expert"], [''] * len(row)
            if exp == "🏆 Team AVG": styles = ['background-color: #cbd5e1; font-weight: 800; color: #0f172a'] * len(row)
            elif "🥇" in exp: styles = ['background-color: #fef08a; color: #854d0e; font-weight: 800'] * len(row)
            elif "🥈" in exp: styles = ['background-color: #e2e8f0; color: #334155; font-weight: 800'] * len(row)
            elif "🥉" in exp: styles = ['background-color: #ffedd5; color: #9a3412; font-weight: 800'] * len(row)
            elif exp == aname and not is_admin(): styles = ['background-color: #dbeafe; color: #1e40af; font-weight: 800'] * len(row)
            if "Prospected Incentive" in row.index and exp != "🏆 Team AVG":
                idx, inc = row.index.get_loc("Prospected Incentive"), str(row["Prospected Incentive"])
                styles[idx] += '; background-color: #fef2f2; color: #dc2626; font-weight: 900; border: 2px solid #fca5a5;' if inc != "3,100 EGP" else '; background-color: #f0fdf4; color: #16a34a; font-weight: 900; border: 2px solid #86efac;'
            return styles

        display_df = display_df[["Expert", "Rank", "Working Days", "Tickets Count", "JHAH Requests", "Support Requests", "Cases/Day", "% Achievement from Target", "AFR", "Service Time", "Service Quality", "Prospected Incentive"]]
        styled_df = display_df.style.apply(style_performers, axis=1).set_properties(**{'text-align': 'center'}).set_properties(subset=['Expert'], **{'font-weight': '900', 'color': '#0f172a'})
        
        try: 
            html_table = styled_df.hide(axis="index").to_html()
        except Exception: 
            try: 
                html_table = styled_df.hide_index().to_html()
            except Exception: 
                html_table = styled_df.to_html()
            
        st.markdown("### 📅 Schedule & Leaves Summary")
        rk1, rk2, rk3, rk4, rk5 = st.columns(5)
        rk1.markdown(kpi_colored("Working Days (Shifts)", int(pd.to_numeric(kpi_scope_df["Working Days"], errors='coerce').fillna(0).sum()), "card-neutral"), unsafe_allow_html=True)
        rk2.markdown(kpi_colored("Off Days", int(pd.to_numeric(kpi_scope_df["Off Days"], errors='coerce').fillna(0).sum()), "card-neutral"), unsafe_allow_html=True)
        rk3.markdown(kpi_colored("Annual Leaves", int(pd.to_numeric(kpi_scope_df["Annual Leaves"], errors='coerce').fillna(0).sum()), "card-neutral"), unsafe_allow_html=True)
        rk4.markdown(kpi_colored("Casual Leaves", int(pd.to_numeric(kpi_scope_df["Casual Leaves"], errors='coerce').fillna(0).sum()), "card-neutral"), unsafe_allow_html=True)
        rk5.markdown(kpi_colored("Sick Leaves", int(pd.to_numeric(kpi_scope_df["Sick Leaves"], errors='coerce').fillna(0).sum()), "card-neutral"), unsafe_allow_html=True)
        
        st.markdown("### 📊 Expert Performance Scorecard Dashboard")
        st.markdown(f'<div class="scorecard-container">{html_table}</div>', unsafe_allow_html=True); st.divider()
        
        st.markdown("### 🔍 Quality Issues Log (Current Period)")
        if not df_q_filtered.empty:
            q_view = df_q_filtered[df_q_filtered['Display_Expert'] == aname].copy() if is_exp else (df_q_filtered[df_q_filtered['Display_Expert'].isin([x for x in OFFICIAL_EXPERTS if x in sel_agents_t2])] if sel_agents_t2 else df_q_filtered.copy())
            if q_view.empty: st.success("🎉 No quality issues recorded for the selected expert(s) in this period!")
            else:
                disp_q = q_view[['Date', 'Expert Name', 'Ticket ID', 'Severity', 'Reason', 'Deduction']].copy().rename(columns={'Deduction': 'Deduction (%)'})
                disp_q['Deduction (%)'] = disp_q['Deduction (%)'].apply(lambda x: f"-{x}%")
                disp_q['Severity'] = disp_q['Severity'].apply(lambda s: {'critical':'🔴 Critical','major':'🟠 Major','medium':'🟠 Medium','minor':'🔵 Minor'}.get(str(s).strip().lower(), s))
                def style_q(row):
                    sev, styles = str(row['Severity']).lower(), ['background-color: #ffffff; color: #1e293b; font-weight: 600'] * len(row)
                    for i, col in enumerate(row.index):
                        if col in ['Severity', 'Deduction (%)']: 
                            if 'critical' in sev: styles[i] = 'background-color: #fee2e2; color: #b91c1c; font-weight: 900;'
                            elif 'major' in sev: styles[i] = 'background-color: #ffedd5; color: #c2410c; font-weight: 900;'
                            elif 'medium' in sev: styles[i] = 'background-color: #ffedd5; color: #c2410c; font-weight: 900;'
                            elif 'minor' in sev: styles[i] = 'background-color: #dbeafe; color: #1d4ed8; font-weight: 900;'
                    return styles
                
                styled_q = disp_q.style.apply(style_q, axis=1).set_properties(**{"text-align": "center"})
                try: 
                    html_q_table = styled_q.hide(axis="index").to_html()
                except Exception: 
                    try: 
                        html_q_table = styled_q.hide_index().to_html()
                    except Exception: 
                        html_q_table = styled_q.to_html()
                    
                st.markdown(f'<div class="scorecard-container">{html_q_table}</div>', unsafe_allow_html=True)
        else: st.info("No quality issues logged in the system for this specific period.")
        
        if is_admin():
            st.divider(); st.markdown("#### 📥 Export Team Performance Report")
            csv_sc = display_df.copy(); csv_sc["Expert"] = csv_sc["Expert"].apply(lambda x: re.sub(r'^[🥇🥈🥉]\s*', '', str(x)))
            st.download_button("📥 Download Team Scorecard (CSV)", csv_sc.to_csv(index=False).encode('utf-8-sig'), f"Team_Scorecard_{d_from}_to_{d_to}.csv", "text/csv", use_container_width=True)
            
            st.divider(); st.markdown("#### ✉️ Performance Review Emails")
            
            sel_email_agent = st.selectbox("Select Agent for Email Draft", [x for x in sc_final["Expert"] if "🏆 Team AVG" not in x])
            if sel_email_agent:
                arow, trow = sc_final[sc_final["Expert"] == sel_email_agent].iloc[0], sc_final[sc_final["Expert"] == "🏆 Team AVG"].iloc[0]
                def sfl(v):
                    try: return float(str(v).replace('%','').replace(',',''))
                    except: return 0.0
                def t2s(t_str):
                    try:
                        p = str(t_str).split(':')
                        return int(p[0])*3600 + int(p[1])*60 + int(p[2])
                    except: return 0
                
                ach_v, q_v = sfl(arow["% Achievement from Target"]), sfl(arow["Service Quality"])
                t_ach_v, t_q_v = sfl(trow["% Achievement from Target"]), sfl(trow["Service Quality"])
                a_tot = int(sfl(arow['Tickets Count']) + sfl(arow['JHAH Requests']) + sfl(arow['Support Requests']))
                t_tot = int(sfl(trow['Tickets Count']) + sfl(trow['JHAH Requests']) + sfl(trow['Support Requests']))
                
                a_afr_sec, t_afr_sec = t2s(arow['AFR']), t2s(trow['AFR'])
                a_st_sec, t_st_sec = t2s(arow['Service Time']), t2s(trow['Service Time'])
                
                c_cases = "#dc2626" if sfl(arow['Cases/Day']) < sfl(trow['Cases/Day']) else "#16a34a"
                c_ach   = "#dc2626" if ach_v < 100 else "#16a34a"
                c_qual  = "#dc2626" if q_v < t_q_v else "#16a34a"
                c_afr   = "#dc2626" if a_afr_sec > t_afr_sec else "#16a34a"
                c_st    = "#dc2626" if a_st_sec > t_st_sec else "#16a34a"
                
                enc_msg = ""
                if ach_v >= 100:
                    perf_w, tgt_m = "outstanding", f"You successfully exceeded the daily target with a brilliant **{arow['% Achievement from Target']}** achievement rate! Keep up the great momentum."
                elif ach_v >= 80:
                    perf_w, tgt_m = "solid", f"You reached a solid **{arow['% Achievement from Target']}** of the daily target. You are very close to hitting the green zone!"
                else:
                    perf_w, tgt_m = "developing", f"You achieved **{arow['% Achievement from Target']}** of the target. We noticed a drop this period, but we know your true capabilities are much higher."
                    enc_msg = "<br><div style='background-color:#eff6ff; padding:15px; border-left:4px solid #3b82f6; margin-top:15px; border-radius:4px;'><b>💡 A Special Note for You:</b><br>Every great expert faces challenging periods. We believe in your potential and are fully here to support you in overcoming these obstacles. Let's work together to boost your numbers and bring out the best in you next week. You've absolutely got this! 💪</div>"

                q_msg = f"Your service quality is top-tier at **{arow['Service Quality']}**." if q_v >= 95 else (f"Your service quality is strong at **{arow['Service Quality']}**." if q_v >= 85 else f"Your service quality sits at **{arow['Service Quality']}**. Let's focus on accuracy and double-check our steps to ensure excellence.")
                if q_v < 90 and not enc_msg:
                    enc_msg = "<br><div style='background-color:#eff6ff; padding:15px; border-left:4px solid #3b82f6; margin-top:15px; border-radius:4px;'><b>💡 A Special Note for You:</b><br>Quality is our top priority, and we noticed some challenges this period. Don't worry, we are here to support you! Let's take a moment to review the quality standards and work together to get your score back to 100%. You can do it! 💪</div>"
                
                c_name = re.sub(r'^[🥇🥈🥉]\s*', '', str(sel_email_agent))
                
                email_html_table = f"""
                <table style="width:100%; border-collapse: collapse; margin: 20px 0; font-family: Arial, sans-serif; font-size: 13px; border: 1px solid #e2e8f0;">
                    <thead>
                        <tr style="background-color: #f8fafc; color: #0f172a; text-align: center; border-bottom: 2px solid #cbd5e1;">
                            <th style="padding: 10px; border-right: 1px solid #e2e8f0; text-align: left;">Metric</th>
                            <th style="padding: 10px; border-right: 1px solid #e2e8f0;">Total Cases</th>
                            <th style="padding: 10px; border-right: 1px solid #e2e8f0;">Cases/Day</th>
                            <th style="padding: 10px; border-right: 1px solid #e2e8f0;">Achievement</th>
                            <th style="padding: 10px; border-right: 1px solid #e2e8f0;">Quality</th>
                            <th style="padding: 10px; border-right: 1px solid #e2e8f0;">AFR</th>
                            <th style="padding: 10px; border-right: 1px solid #e2e8f0;">Service Time</th>
                            <th style="padding: 10px;">Incentive</th>
                        </tr>
                    </thead>
                    <tbody>
                        <tr style="border-bottom: 1px solid #e2e8f0;">
                            <td style="padding: 10px; border-right: 1px solid #e2e8f0; font-weight: bold; color: #334155; white-space: nowrap;">Your Score 👤</td>
                            <td style="padding: 10px; border-right: 1px solid #e2e8f0; text-align: center; font-weight: bold;">{a_tot}</td>
                            <td style="padding: 10px; border-right: 1px solid #e2e8f0; text-align: center; color: {c_cases}; font-weight: bold;">{arow['Cases/Day']}</td>
                            <td style="padding: 10px; border-right: 1px solid #e2e8f0; text-align: center; color: {c_ach}; font-weight: bold;">{arow['% Achievement from Target']}</td>
                            <td style="padding: 10px; border-right: 1px solid #e2e8f0; text-align: center; color: {c_qual}; font-weight: bold;">{arow['Service Quality']}</td>
                            <td style="padding: 10px; border-right: 1px solid #e2e8f0; text-align: center; color: {c_afr}; font-weight: bold;">{arow['AFR']}</td>
                            <td style="padding: 10px; border-right: 1px solid #e2e8f0; text-align: center; color: {c_st}; font-weight: bold;">{arow['Service Time']}</td>
                            <td style="padding: 10px; text-align: center; color: #16a34a; font-weight: bold;">{arow['Prospected Incentive']}</td>
                        </tr>
                        <tr style="background-color: #fcfcfc;">
                            <td style="padding: 10px; border-right: 1px solid #e2e8f0; font-weight: bold; color: #334155; white-space: nowrap;">Team Avg 🏆</td>
                            <td style="padding: 10px; border-right: 1px solid #e2e8f0; text-align: center; color: #475569;">{t_tot}</td>
                            <td style="padding: 10px; border-right: 1px solid #e2e8f0; text-align: center; color: #475569;">{trow['Cases/Day']}</td>
                            <td style="padding: 10px; border-right: 1px solid #e2e8f0; text-align: center; color: #475569;">{trow['% Achievement from Target']}</td>
                            <td style="padding: 10px; border-right: 1px solid #e2e8f0; text-align: center; color: #475569;">{trow['Service Quality']}</td>
                            <td style="padding: 10px; border-right: 1px solid #e2e8f0; text-align: center; color: #475569;">{trow['AFR']}</td>
                            <td style="padding: 10px; border-right: 1px solid #e2e8f0; text-align: center; color: #475569;">{trow['Service Time']}</td>
                            <td style="padding: 10px; text-align: center; font-weight: bold; color: #475569;">3,100 EGP</td>
                        </tr>
                    </tbody>
                </table>
                """
                
                st.markdown("##### 📝 Email Preview (Highlight & Copy directly from here!)")
                st.markdown(f"<div style='background:#ffffff; padding:2rem; border-radius:12px; border:1px solid #cbd5e1; font-size:1.1rem; color:#334155;'>Dear **{c_name}**,<br><br>As we review the performance for the period from **{d_from}** to **{d_to}**, I wanted to share your metrics and highlight your **{perf_w}** contributions.<br>{email_html_table}<b>🎯 Targets & Quality:</b><br>{tgt_m}<br>{q_msg}{enc_msg}<br><br>Thank you for your hard work!<br><br>Best regards,<br><b>Mohammed Shehta</b><br>Team Leader</div>", unsafe_allow_html=True)
                
                st.markdown(f'<a href="https://mail.google.com/mail/?view=cm&fs=1&to=&su={urllib.parse.quote(f"Your Performance Review ({d_from} to {d_to}) - {c_name}")}" target="_blank" style="display:block; padding:0.8rem 1.2rem; background-color:#2563eb; color:white; text-decoration:none; border-radius:8px; font-weight:800; font-size:1.15rem; width:100%; text-align:center; margin-top: 10px; box-shadow: 0 4px 6px rgba(37,99,235, 0.3);">🌐 Open Draft in Gmail</a>', unsafe_allow_html=True)

# ── TAB 3 — Manual Overrides (Admin Only) ─────────────────────────────────────────
if is_admin():
    with tabs[2]:
        st.markdown("### 🎯 Global Settings (This Period)")
        with st.form("global_target_form"):
            new_target = st.number_input("Daily Target (Cases/Day)", value=int(global_target), step=1, min_value=0)
            if st.form_submit_button("💾 Save Global Parameters", use_container_width=True):
                if PERIOD_KEY not in overrides(): overrides()[PERIOD_KEY] = {}
                overrides()[PERIOD_KEY]["GLOBAL_TARGET"] = new_target; _save_store(); st.success("✅ Saved!"); st.rerun()

        st.markdown("### 🤫 Secret Agent Ticket Adjustments")
        st.info("💡 Any number added here will be invisibly added to the agent's total tickets for this period. Experts will NOT see the adjustment itself, only the final increased/decreased total.")
        
        if "agent_adjustments" not in overrides().get(PERIOD_KEY, {}):
            if PERIOD_KEY not in overrides(): overrides()[PERIOD_KEY] = {}
            overrides()[PERIOD_KEY]["agent_adjustments"] = {}
            
        current_adjs = overrides()[PERIOD_KEY].get("agent_adjustments", {})
        
        with st.form("secret_adj_form"):
            adj_agent = st.selectbox("Select Agent", OFFICIAL_EXPERTS)
            adj_val = st.number_input("Adjustment Value (can be negative)", value=int(current_adjs.get(adj_agent, 0)), step=1)
            
            if st.form_submit_button("💾 Save Secret Adjustment", use_container_width=True):
                if PERIOD_KEY not in overrides(): overrides()[PERIOD_KEY] = {}
                if "agent_adjustments" not in overrides()[PERIOD_KEY]: overrides()[PERIOD_KEY]["agent_adjustments"] = {}
                
                overrides()[PERIOD_KEY]["agent_adjustments"][adj_agent] = adj_val
                _save_store()
                st.success(f"✅ Secret adjustment of {adj_val} applied to {adj_agent}!")
                st.rerun()
        
        if current_adjs:
            st.markdown("**Current Active Adjustments for this period:**")
            st.json(current_adjs)

# ── TAB 4 — Team Profiles (Admin & Expert) ────────────────────────────────────────
if len(tabs) > 1 and (is_admin() or st.session_state.role == "expert"):
    with tabs[-1]:
        st.markdown("### 🧑‍🤝‍🧑 Approvals Team Profiles")
        st.info("💡 تعرف على زملائك في الفريق، تاريخ انضمامهم، وخبراتهم!")
        
        team_members = [u for uid, u in users().items() if u["role"] in ["expert", "admin", "supervisor"]]
        cols = st.columns(3)
        for i, member in enumerate(team_members):
            photo_url = member.get("photo", "") or "https://cdn-icons-png.flaticon.com/512/3135/3135715.png" 
            r_badge = "👑 ADMIN" if member['role'] == 'admin' else ("👁️ SUPERVISOR" if member['role'] == 'supervisor' else "⭐ EXPERT")
            grad, join_cc, join_team = member.get('grad_year') or 'N/A', member.get('join_cc') or 'N/A', member.get('join_team') or 'N/A'
            bio = member.get('bio') or 'لا توجد نبذة شخصية حتى الآن.'
            
            card_html = f"""
            <div class="profile-card">
                <img src="{photo_url}" class="profile-img" onerror="this.src='https://cdn-icons-png.flaticon.com/512/3135/3135715.png'">
                <div class="profile-name">{member.get('display_name', 'Unknown')}</div>
                <div class="profile-role">{r_badge}</div>
                <div class="profile-detail"><span>🎓 Graduation</span> <b>{grad}</b></div>
                <div class="profile-detail"><span>🏢 Joined CC</span> <b>{join_cc}</b></div>
                <div class="profile-detail"><span>🚀 Joined Team</span> <b>{join_team}</b></div>
                <div class="profile-bio">"{bio}"</div>
            </div>
            """
            cols[i % 3].markdown(card_html, unsafe_allow_html=True)

st.info(f"⏱️ Operational Sync Status: Metrics loaded completely across {len(df)} synced records.")
# --- END OF SCRIPT ---
